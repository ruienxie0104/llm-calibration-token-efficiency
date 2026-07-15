#!/usr/bin/env python3
"""
Research Proposal Slides: LLM Calibration x Token Efficiency
For advisor meeting 2026-07-04
Style: No emoji, light color scheme, clean typography (matching Ryan's preference)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
TEAL = RGBColor(0x2E, 0x8B, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_BLUE = RGBColor(0xE8, 0xEE, 0xF4)
LIGHT_TEAL = RGBColor(0xE8, 0xF4, 0xF4)
ALT_ROW = RGBColor(0xF0, 0xF4, 0xF8)
MEDIUM_GRAY = RGBColor(0x66, 0x6E, 0x7A)
DARK_TEXT = RGBColor(0x2D, 0x3A, 0x4A)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED = RGBColor(0xFD, 0xED, 0xED)
GREEN_ACCENT = RGBColor(0x27, 0x84, 0x3C)
LIGHT_GREEN = RGBColor(0xE8, 0xF4, 0xE8)
AMBER = RGBColor(0xB8, 0x86, 0x0B)
LIGHT_AMBER = RGBColor(0xFB, 0xF3, 0xE8)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# ============================================================
# Helper functions
# ============================================================

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    return add_shape(slide, left, top, width, height, fill_color, line_color, line_width, MSO_SHAPE.ROUNDED_RECTANGLE)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14, color=DARK_TEXT, spacing=Pt(6)):
    """Add multiple paragraphs in one textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get('size', font_size))
        run.font.bold = opts.get('bold', False)
        run.font.color.rgb = opts.get('color', color)
        run.font.name = 'Calibri'
    return txBox

def add_footer(slide, page_num, total=22):
    add_textbox(slide, Emu(11155680), Emu(6400800), Emu(914400), Emu(365760),
                f'{page_num} / {total}', font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    add_shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(54864), fill_color=TEAL)

def add_bottom_bar(slide):
    add_shape(slide, Emu(0), Emu(6537960), SLIDE_W, Emu(320040), fill_color=DARK_BLUE)

def add_title_bar(slide, title_text, subtitle_text=None):
    add_shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(1188720), fill_color=DARK_BLUE)
    add_textbox(slide, Emu(548640), Emu(137160), Emu(10972800), Emu(640080),
                title_text, font_size=32, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Emu(548640), Emu(685800), Emu(10972800), Emu(457200),
                    subtitle_text, font_size=16, color=RGBColor(0xD0, 0xEC, 0xEC))

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def add_left_border_card(slide, left, top, width, height, border_color=TEAL, fill_color=None):
    if fill_color:
        add_shape(slide, left, top, width, height, fill_color=fill_color)
    add_shape(slide, left, top, Emu(54864), height, fill_color=border_color)

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK_TEXT, bullet='bullet'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f'  {item}'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return txBox

def add_table_simple(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a simple styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.LEFT
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1B3A5C'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK_TEXT
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.LEFT
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            if row_idx % 2 == 0:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'F0F4F8'})
            else:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF'})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

    return table_shape

def add_arrow(slide, left, top, width, height):
    """Add a right-pointing arrow."""
    return add_shape(slide, left, top, width, height, fill_color=TEAL, shape_type=MSO_SHAPE.RIGHT_ARROW)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)

# Accent lines
add_shape(slide, Emu(3657600), Emu(2286000), Emu(4876495), Emu(54864), fill_color=TEAL)
add_shape(slide, Emu(3657600), Emu(4389120), Emu(4876495), Emu(54864), fill_color=TEAL)

add_textbox(slide, Emu(1371600), Emu(2468880), Emu(9448495), Emu(1097280),
            'LLM Self-Assessment Calibration\nx Token Efficiency', font_size=40, bold=True, color=DARK_BLUE)

add_textbox(slide, Emu(1828800), Emu(3749039), Emu(8534095), Emu(457200),
            'Research Proposal for Advisor Meeting', font_size=18, color=MEDIUM_GRAY)

add_textbox(slide, Emu(1828800), Emu(4754880), Emu(8534095), Emu(731520),
            'Can calibrated self-evaluation predict and optimize token allocation?', font_size=20, color=TEAL)

add_textbox(slide, Emu(2743200), Emu(5669280), Emu(6705295), Emu(457200),
            'Ryan  -  2026-07-04', font_size=16, color=MEDIUM_GRAY)

add_footer(slide, 1)
add_notes(slide, "老師好，今天要跟您討論的是延續學姐論文的後續研究方向。學姐證明了 IRT 難度信號能改善 LLM 的信心校準，也提到 reliability 跟 inference cost 有關聯，但沒有深入。我的研究想填補這個缺口：驗證校準好的模型能不能用更少 token 達到同樣品質。接下來我會從學姐的框架開始，帶到文獻調查結果、研究問題、實驗設計，最後有幾個問題想跟老師討論。")


# ============================================================
# SLIDE 2: Outline
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_top_bar(slide)
add_textbox(slide, Emu(548640), Emu(137160), Emu(10972800), Emu(640080),
            'Outline', font_size=32, bold=True, color=DARK_BLUE)
add_bottom_bar(slide)

outline_items = [
    '1. Background: 學姐論文回顧與缺口',
    '2. Research Direction: Token 最小使用',
    '3. Literature Survey: 16 组 arXiv 搜索结果',
    '4. Four Core Gaps',
    '5. Research Questions & Hypotheses',
    '6. Paper Story & Contributions',
    '7. Experiment Design (Phase 1-3)',
    '8. Integration with 學姐 Framework',
    '9. Risk Assessment',
    '10. Timeline & Questions for Advisor',
]

y_start = 1463040
y_step = 475488

for i, item in enumerate(outline_items):
    y = y_start + y_step * i
    add_shape(slide, Emu(1371600), Emu(y), Emu(411480), Emu(411480),
              fill_color=TEAL, shape_type=MSO_SHAPE.OVAL)
    # Number text
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = str(i+1)
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    # Item text
    add_textbox(slide, Emu(2011680), Emu(y + 45720), Emu(8229600), Emu(365760),
                item[3:] if item[0].isdigit() else item, font_size=18, color=DARK_TEXT)

add_footer(slide, 2)
add_notes(slide, "今天的討論分十個部分。先回顧學姐的論文和缺口，再講老師建議的 token 方向，然後是我做的深度文獻調查、四個核心缺口、研究問題、貢獻設計、實驗規劃，最後是風險評估和給老師的問題。")


# ============================================================
# SLIDE 3: 學姐論文回顧
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Background: 學姐論文回顧', 'Chen et al., IEEE IRI 2026 - LCAE Framework')
add_bottom_bar(slide)

# Three stages
stages = [
    ('Stage 1: 能力估計', 'Rasch Model\n模型能力 + 題目難度\n同一把尺\n客觀算答錯機率', DARK_BLUE, LIGHT_BLUE),
    ('Stage 2: 自我評估', '四種情境：\nQOQ (不给訊號)\nIDS (給難度訊號)\nDPR (給能力位置)\nCombined (全給)', TEAL, LIGHT_TEAL),
    ('Stage 3: 模型選擇', 'LCAE 指標\n比較「模型自估」\nvs「IRT 客觀估」\n一致性 + 成本效益', RGBColor(0x6A, 0x4C, 0x93), RGBColor(0xF0, 0xEC, 0xF8)),
]

card_w = Emu(3403520)
card_h = Emu(3200400)
card_y = Emu(1463040)
gap = Emu(365760)

for i, (title, desc, border, fill) in enumerate(stages):
    x = Emu(457200 + (3403520 + 365760) * i)
    add_left_border_card(slide, x, card_y, card_w, card_h, border_color=border, fill_color=fill)
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 91440), Emu(3000000), Emu(457200),
                title, font_size=16, bold=True, color=border)
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 548640), Emu(3000000), Emu(2400000),
                desc, font_size=13, color=DARK_TEXT)

# Key findings
add_shape(slide, Emu(457200), Emu(4800600), Emu(11277600), Emu(1463040), fill_color=LIGHT_GRAY)
add_textbox(slide, Emu(548640), Emu(4846320), Emu(1097280), Emu(365760),
            'Key Findings', font_size=16, bold=True, color=DARK_BLUE)
findings_text = (
    "能力強 != 自評準 (GPT-5 最強但 LCAE 不是最好; Llama 3 70B 自評最準)    "
    "IDS 最有效 (給難度訊號是改善校準最有效的元素)    "
    "校準不傷能力    Cost 提到但未深入"
)
add_textbox(slide, Emu(548640), Emu(5212320), Emu(10972800), Emu(822960),
            findings_text, font_size=13, color=DARK_TEXT)

add_footer(slide, 3)
add_notes(slide, "學姐的論文分三個階段。第一階段用 Rasch Model 把模型能力和題目難度放在同一把尺上，可以客觀算出答錯機率。第二階段設計了四種自我評估情境，從完全不給訊號到全給，讓模型自估錯誤機率。第三階段用 LCAE 指標比較模型自估和 IRT 客觀估的一致性，搭配成本效益做模型選擇。三個關鍵發現：能力強不等於自評準、IDS 最有效、校準不傷能力。但 cost 只有粗估，token 議題完全沒驗證。")


# ============================================================
# SLIDE 4: 學姐論文的四個缺口
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, '學姐論文的四個缺口', 'Where this research can extend')
add_bottom_bar(slide)

gaps = [
    ('1', '只有 Rasch Model', '未計算鑑別度\n可引入 2PL/3PL', RED_ACCENT, LIGHT_RED),
    ('2', '難度訊號偏差', 'IDS 的難度訊號來自\n同組受評模型，有循環依賴\n可用外部獨立模型', AMBER, LIGHT_AMBER),
    ('3', '成本估計簡化', '未應用於實際場景\n只有粗略 cost 估計', AMBER, LIGHT_AMBER),
    ('4', 'Token 未驗證', 'token 使用量與答題品質\n自評準確度的關係\n完全未驗證', RED_ACCENT, LIGHT_RED),
]

card_w = Emu(2700000)
card_h = Emu(3200400)
card_y = Emu(1463040)

for i, (num, title, desc, border, fill) in enumerate(gaps):
    x = Emu(457200 + (2700000 + 228600) * i)
    add_left_border_card(slide, x, card_y, card_w, card_h, border_color=border, fill_color=fill)
    # Number badge
    add_shape(slide, Emu(x.emu + 137160), Emu(card_y.emu + 91440), Emu(411480), Emu(411480),
              fill_color=border, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    # Title
    add_textbox(slide, Emu(x.emu + 640080), Emu(card_y.emu + 137160), Emu(1900000), Emu(457200),
                title, font_size=16, bold=True, color=border)
    # Description
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 731520), Emu(2400000), Emu(2200000),
                desc, font_size=13, color=DARK_TEXT)

# Teacher's suggestion box
add_shape(slide, Emu(457200), Emu(4983480), Emu(11277600), Emu(1280160), fill_color=LIGHT_TEAL)
add_shape(slide, Emu(457200), Emu(4983480), Emu(54864), Emu(1280160), fill_color=TEAL)
add_textbox(slide, Emu(548640), Emu(5029200), Emu(10972800), Emu(365760),
            "Teacher's Suggestion for Ryan", font_size=16, bold=True, color=TEAL)
add_textbox(slide, Emu(548640), Emu(5392680), Emu(10972800), Emu(822960),
            "Extend the token minimization issue from 學姐's framework. Design simple vs complex task scenarios, "
            "validate the relationship between token usage and answer quality / self-assessment accuracy. "
            "Goal: achieve equal quality with fewer tokens.", font_size=14, color=DARK_TEXT)

add_footer(slide, 4)
add_notes(slide, "老師之前指出了學姐論文的四個缺口。前三個是用更精細的 IRT 模型、解決難度訊號偏差、跟把成本估計做得更精確。但第四個最重要，也是老師給我的方向：token 使用量跟答題品質和自評準確度的關係完全沒驗證。老師建議我延續這個方向，設計簡單跟複雜任務的場景，看 token 使用量和品質的關係，目標是用更少 token 達到同等品質。")


# ============================================================
# SLIDE 5: 研究核心直覺
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Research Direction: Core Intuition', 'Three questions 學姐 did not answer')
add_bottom_bar(slide)

questions = [
    ("Q1", "改善信心校準之後，模型能不能更有效率地使用 token？", DARK_BLUE, LIGHT_BLUE),
    ("Q2", "如果模型「知道自己會不會」，能不能用更少 token 達到同樣品質？", TEAL, LIGHT_TEAL),
    ("Q3", "自我評估準確度 (LCAE) 和 token 使用效率之間有沒有關係？", RGBColor(0x6A, 0x4C, 0x93), RGBColor(0xF0, 0xEC, 0xF8)),
]

card_y = Emu(1463040)
card_h = Emu(1371600)
card_w = Emu(11277600)

for i, (label, text, border, fill) in enumerate(questions):
    y = card_y.emu + (card_h.emu + 228600) * i
    add_left_border_card(slide, Emu(457200), Emu(y), card_w, card_h, border_color=border, fill_color=fill)
    # Label badge
    add_shape(slide, Emu(640080), Emu(y + 137160), Emu(548640), Emu(548640),
              fill_color=border, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    # Question text
    add_textbox(slide, Emu(1371600), Emu(y + 228600), Emu(10000000), Emu(822960),
                text, font_size=20, bold=True, color=DARK_TEXT)

# Bottom summary
add_shape(slide, Emu(457200), Emu(5669280), Emu(11277600), Emu(548640), fill_color=DARK_BLUE)
add_textbox(slide, Emu(548640), Emu(5715240), Emu(10972800), Emu(365760),
            "From correlation (學姐 mentioned) -> causation (this research validates) -> optimization (token allocation strategy)",
            font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 5)
add_notes(slide, "學姐證明了給難度信號能改善校準，但留下了三個沒回答的問題。第一，改善校準之後模型能不能更有效率用 token？第二，如果模型知道自己會不會，能不能少用 token？第三，自我評估準確度和 token 效率有沒有關係？這三個問題就是我的研究主軸。底部的橘色條是整個研究的邏輯：從學姐觀察到的關聯，到本研究要驗證的因果，再到最終的 token 分配優化。")


# ============================================================
# SLIDE 6: 文獻調查規模
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Literature Survey: 16 Systematic arXiv Searches', 'Covering self-evaluation, calibration, token efficiency, IRT, and their intersections')
add_bottom_bar(slide)

# Big number
add_textbox(slide, Emu(457200), Emu(1463040), Emu(1828800), Emu(1097280),
            '16', font_size=80, bold=True, color=TEAL)
add_textbox(slide, Emu(2286000), Emu(1645920), Emu(2743200), Emu(457200),
            'systematic searches', font_size=20, color=DARK_TEXT)
add_textbox(slide, Emu(2286000), Emu(2194560), Emu(2743200), Emu(365760),
            'on arXiv (cs.AI, cs.CL, cs.LG)', font_size=14, color=MEDIUM_GRAY)

# 6 zero-result queries
add_textbox(slide, Emu(457200), Emu(2743200), Emu(5486400), Emu(457200),
            '6 precise intersection queries returned 0 results:', font_size=18, bold=True, color=RED_ACCENT)

zero_queries = [
    '"difficulty" x "token allocation" x "LLM"',
    '"adaptive computation" x "self-evaluation" x "LLM"',
    '"IRT" x "LLM evaluation" x "cost"',
    '"self-assessment" x "reasoning cost"',
    '"compute allocation" x "self-knowledge"',
    '"self-calibration" x "reasoning" x "efficiency"',
]

for i, q in enumerate(zero_queries):
    col = i % 2
    row = i // 2
    x = Emu(457200 + col * 5486400)
    y = Emu(3200400 + row * 457200)
    add_shape(slide, x, y, Emu(320040), Emu(320040), fill_color=RED_ACCENT, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = '0'
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_textbox(slide, Emu(x.emu + 365760), Emu(y.emu + 45720), Emu(4800000), Emu(320040),
                q, font_size=13, color=DARK_TEXT)

# Right side: 5 clusters
add_textbox(slide, Emu(6400800), Emu(1463040), Emu(5486400), Emu(457200),
            '5 existing literature clusters', font_size=18, bold=True, color=TEAL)

clusters = [
    ('A', 'Difficulty-aware Token Budget', '~10 papers', DARK_BLUE),
    ('B', 'LLM Confidence Calibration', '~8 papers', TEAL),
    ('C', 'Test-Time Compute Scaling', '~9 papers', RGBColor(0x6A, 0x4C, 0x93)),
    ('D', 'IRT x LLM Evaluation', 'only 學姐', AMBER),
    ('E', 'Overthinking / Redundancy', 'several', MEDIUM_GRAY),
]

for i, (label, name, count, color) in enumerate(clusters):
    y = Emu(2011680 + 548640 * i)
    add_shape(slide, Emu(6400800), y, Emu(411480), Emu(411480), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_textbox(slide, Emu(6943680), Emu(y.emu + 45720), Emu(3200000), Emu(365760),
                name, font_size=15, bold=True, color=DARK_TEXT)
    add_textbox(slide, Emu(10000000), Emu(y.emu + 45720), Emu(1645920), Emu(365760),
                count, font_size=13, color=MEDIUM_GRAY)

# Key insight
add_shape(slide, Emu(457200), Emu(5486400), Emu(11277600), Emu(640080), fill_color=LIGHT_AMBER)
add_shape(slide, Emu(457200), Emu(5486400), Emu(54864), Emu(640080), fill_color=AMBER)
add_textbox(slide, Emu(548640), Emu(5531880), Emu(10972800), Emu(457200),
            "Key insight: These 5 clusters do not intersect with each other. Nobody combines self-assessment + IRT + token prediction.",
            font_size=14, bold=True, color=DARK_TEXT)

add_footer(slide, 6)
add_notes(slide, "我做了 16 組系統性的 arXiv 搜尋，覆蓋自我評估、校準、token 效率、IRT 等關鍵字。最關鍵的發現是：6 組精準交集查詢全部返回零結果。比如 difficulty 乘 token allocation 乘 LLM、IRT 乘 LLM evaluation 乘 cost、self-assessment 乘 reasoning cost，這些都是零。現有文獻分成五個集群，但彼此不交集。集群 A 做難度感知的 token 分配但用外部信號，集群 B 做信心校準但不碰 token，集群 C 做 test-time compute 但用 entropy 或 RL，集群 D 只有學姐一篇用 IRT，集群 E 分析冗餘但不是從自我評估角度。沒有人把這些連起來。")


# ============================================================
# SLIDE 7: 文獻空白地圖
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Literature Gap Map', 'Where the white space is')
add_bottom_bar(slide)

# Axis
add_shape(slide, Emu(5486400), Emu(5765040), Emu(548640), Emu(320040), fill_color=MEDIUM_GRAY)  # x axis arrow
add_shape(slide, Emu(5486400), Emu(1554480), Emu(548640), Emu(4389120), fill_color=MEDIUM_GRAY)  # y axis

add_textbox(slide, Emu(5486400), Emu(6126480), Emu(548640), Emu(320040),
            'External Signal', font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(457200), Emu(1371600), Emu(4572000), Emu(320040),
            'Token Efficiency', font_size=12, color=MEDIUM_GRAY)
add_textbox(slide, Emu(457200), Emu(1695960), Emu(4572000), Emu(320040),
            '(higher = better)', font_size=10, color=MEDIUM_GRAY)
add_textbox(slide, Emu(6217920), Emu(1371600), Emu(4572000), Emu(320040),
            'Self-Assessment', font_size=12, color=MEDIUM_GRAY)

# Existing work (left side - external signals)
existing_left = [
    ('UAB', 1371600, 2743200),
    ('ROI-Reasoning', 1371600, 3200400),
    ('CARD', 1828800, 2468880),
    ('TAB', 1828800, 2926080),
    ('DSC', 2011680, 2286000),
    ('"LLM Already Knows"', 2286000, 2743200),
]

for name, y, x in existing_left:
    add_shape(slide, Emu(x), Emu(y), Emu(1280160), Emu(320040),
              fill_color=LIGHT_BLUE, line_color=DARK_BLUE, line_width=Pt(1))
    add_textbox(slide, Emu(x + 45720), Emu(y + 45720), Emu(1180000), Emu(228600),
                name, font_size=10, color=DARK_BLUE)

# Existing work (right side - self-assessment, but no token)
existing_right = [
    ('學姐 LCAE', 5303520, 2194560),
    ('Kumaran', 5765040, 2468880),
    ('Protocol Sensitivity', 5486400, 2743200),
    ('Score Granularity', 5303520, 2926080),
    ('Query-Level Uncert.', 5765040, 3200400),
]

for name, x, y in existing_right:
    add_shape(slide, Emu(x), Emu(y), Emu(1645920), Emu(320040),
              fill_color=LIGHT_TEAL, line_color=TEAL, line_width=Pt(1))
    add_textbox(slide, Emu(x + 45720), Emu(y + 45720), Emu(1540000), Emu(228600),
                name, font_size=10, color=TEAL)

# GAPS
gaps_info = [
    ('GAP 1 (star)', 'Self-assessment accuracy\nx Token prediction\nNOBODY', 7680960, 2011680, RED_ACCENT, LIGHT_RED),
    ('GAP 4 (star)', 'Calibration improvement\n-> Token savings\ncausal chain unverified', 7680960, 3200400, RED_ACCENT, LIGHT_RED),
    ('GAP 2', 'IRT x Token allocation\nNOBODY', 7132320, 4114800, AMBER, LIGHT_AMBER),
    ('GAP 3', 'Verbalized conf as\ntoken signal\nUNTESTED', 7680960, 4663440, AMBER, LIGHT_AMBER),
]

for label, desc, x, y, border, fill in gaps_info:
    add_left_border_card(slide, Emu(x), Emu(y), Emu(2743200), Emu(731520),
                         border_color=border, fill_color=fill)
    add_textbox(slide, Emu(x + 91440), Emu(y + 45720), Emu(2550000), Emu(640080),
                desc, font_size=11, bold=True, color=border)

# Test-time compute cluster
tt_items = [
    ('Entropy-Gated', 2743200, 3429360),
    ('SLAT', 3200400, 3429360),
    ('TRIAGE', 3657600, 3429360),
    ('RLCM', 2743200, 3885600),
    ('SR2AM', 3200400, 3885600),
]
for name, y, x in tt_items:
    add_shape(slide, Emu(x), Emu(y), Emu(1180000), Emu(274320),
              fill_color=RGBColor(0xF0, 0xEC, 0xF8), line_color=RGBColor(0x6A, 0x4C, 0x93), line_width=Pt(1))
    add_textbox(slide, Emu(x + 45720), Emu(y + 45720), Emu(1080000), Emu(182880),
                name, font_size=10, color=RGBColor(0x6A, 0x4C, 0x93))

add_textbox(slide, Emu(2743200), Emu(4297680), Emu(3000000), Emu(228600),
            'Test-Time Compute Scaling', font_size=12, bold=True, color=RGBColor(0x6A, 0x4C, 0x93))

add_footer(slide, 7)
add_notes(slide, "這張圖是文獻空白地圖。縱軸是 token 效率，橫軸從外部信號到自我評估。左邊是集群 A 的論文，全用外部信號做 token 分配。右邊是集群 B 的論文，做信心校準但不碰 token。中間偏上是 test-time compute 集群。四個紅色和黃色的框就是四個 gap。Gap 1 和 Gap 4 是三星級的主要空白：自我評估準確度能不能預測 token 需求、校準改善能不能導致 token 省下的因果鏈，都沒人做。Gap 2 和 3 是二星級：IRT 乘 token 和 verbalized confidence 做信號，也沒人測過。重點是這五個集群彼此不交集，我們的貢獻就是把它們連起來。")


# ============================================================
# SLIDE 8: 四個核心缺口
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Four Core Gaps', 'Confirmed by 16 systematic arXiv searches')
add_bottom_bar(slide)

gap_data = [
    ('Gap 1', 'star', 'Self-assessment accuracy -> Token prediction', 'NOBODY has studied whether LLM self-assessment quality can predict token requirements', RED_ACCENT, LIGHT_RED),
    ('Gap 2', 'important', 'IRT difficulty signals x Token allocation', 'Nobody uses IRT objective difficulty to guide token distribution', AMBER, LIGHT_AMBER),
    ('Gap 3', 'risk', 'Verbalized confidence as token signal', 'UNTESTED. Kumaran found verbalized conf tracks commitment, not correctness', AMBER, LIGHT_AMBER),
    ('Gap 4', 'star', 'Calibration improvement -> Token savings causal chain', 'Nobody connects calibration improvement to token efficiency gains', RED_ACCENT, LIGHT_RED),
]

card_w = Emu(5400000)
card_h = Emu(2286000)
positions = [(457200, 1463040), (6217920, 1463040), (457200, 3894720), (6217920, 3894720)]

for i, (label, level, title, desc, border, fill) in enumerate(gap_data):
    x, y = positions[i]
    add_left_border_card(slide, Emu(x), Emu(y), Emu(card_w), Emu(card_h), border_color=border, fill_color=fill)
    # Label
    add_textbox(slide, Emu(x + 137160), Emu(y + 91440), Emu(1200000), Emu(365760),
                label, font_size=20, bold=True, color=border)
    # Level badge
    level_text = {'star': 'star star star', 'important': 'star star', 'risk': 'star star'}.get(level, 'star')
    add_textbox(slide, Emu(x + 1645920), Emu(y + 91440), Emu(1645920), Emu(365760),
                level_text, font_size=14, color=AMBER)
    # Title
    add_textbox(slide, Emu(x + 137160), Emu(y + 548640), Emu(5100000), Emu(457200),
                title, font_size=16, bold=True, color=DARK_TEXT)
    # Description
    add_textbox(slide, Emu(x + 137160), Emu(y + 1097280), Emu(5100000), Emu(1000000),
                desc, font_size=13, color=DARK_TEXT)

add_footer(slide, 8)
add_notes(slide, "四個核心缺口。Gap 1 是最主要的：沒有人研究模型自我評估有多準跟它需要多少 token 之間的關係。如果這成立，我們就有了一個 training-free 的 token 分配信號，因為模型是自己知道的。Gap 2 是 IRT 的客觀難度可以用來做 token 分配，也沒人做。Gap 3 是 verbalized confidence 能不能做 token 分配信號，Kumaran 發現它追蹤 commitment 不是 correctness，這是風險也是研究問題。Gap 4 是核心機會：校準改善到 token 省下的因果鏈沒人驗證過。學姐證明了 IDS 改善校準，有人證明了 difficulty-aware allocation 減少 token，但沒有人把這兩件事連起來。")


# ============================================================
# SLIDE 9: 研究問題與假設
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Research Questions & Hypotheses')
add_bottom_bar(slide)

# RQs on left
add_textbox(slide, Emu(457200), Emu(1463040), Emu(5486400), Emu(457200),
            'Research Questions', font_size=20, bold=True, color=DARK_BLUE)

rqs = [
    ('RQ1', 'Can LCAE predict token requirements?', 'If model knows itself, does it naturally use fewer tokens?'),
    ('RQ2', 'Can IRT difficulty signals serve as token allocation signals?', 'IDS improved calibration; can it also improve token distribution?'),
    ('RQ3', 'Does calibration improvement causally lead to token efficiency?', 'Not just correlation: improvement -> prediction -> savings'),
    ('RQ4', 'Verbalized vs log-prob: which is better for token allocation?', 'Kumaran found verbalized tracks commitment, not correctness'),
]

y = 2011680
for label, q, detail in rqs:
    add_shape(slide, Emu(457200), Emu(y), Emu(548640), Emu(548640),
              fill_color=TEAL, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_textbox(slide, Emu(1180000), Emu(y + 22860), Emu(4800000), Emu(365760),
                q, font_size=14, bold=True, color=DARK_TEXT)
    add_textbox(slide, Emu(1180000), Emu(y + 320040), Emu(4800000), Emu(228600),
                detail, font_size=12, color=MEDIUM_GRAY)
    y += 1097280

# Hypotheses on right
add_textbox(slide, Emu(6400800), Emu(1463040), Emu(5486400), Emu(457200),
            'Hypotheses', font_size=20, bold=True, color=DARK_BLUE)

hyps = [
    ('H1', 'LCAE positively correlates with token efficiency', 'Medium-High confidence', GREEN_ACCENT),
    ('H2', 'IDS difficulty signals improve token allocation efficiency', 'Medium-High confidence', GREEN_ACCENT),
    ('H3', 'Calibration improvement -> token efficiency causal chain holds', 'Medium (null result is also a finding)', AMBER),
    ('H4', 'Log-prob confidence > verbalized confidence for token allocation', 'Medium (Kumaran evidence supports)', AMBER),
]

y = 2011680
for label, h, detail, color in hyps:
    add_shape(slide, Emu(6400800), Emu(y), Emu(548640), Emu(548640),
              fill_color=color, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    add_textbox(slide, Emu(7132320), Emu(y + 22860), Emu(4800000), Emu(365760),
                h, font_size=14, bold=True, color=DARK_TEXT)
    add_textbox(slide, Emu(7132320), Emu(y + 320040), Emu(4800000), Emu(228600),
                detail, font_size=12, color=MEDIUM_GRAY)
    y += 1097280

add_footer(slide, 9)
add_notes(slide, "四個研究問題。RQ1 是 LCAE 能不能預測 token 需求，RQ2 是 IDS 能不能做 token 分配信號，RQ3 是校準改善到 token 省的因果鏈成不成立，RQ4 是 verbalized 和 log-prob 哪個更適合。對應四個假設。H1 和 H2 信心中高，H3 是最關鍵的但只有中等信心因為因果鏈可能不成立，H4 有 Kumaran 的證據支持。特別強調 H3：如果不成立，這本身也是一個重要發現，打破了直覺假設。")


# ============================================================
# SLIDE 10: 論文 Story
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Paper Story', 'The causal chain we want to validate')
add_bottom_bar(slide)

# Chain visualization
chain_items = [
    ('IRT Difficulty\nSignals', '學姐已證', DARK_BLUE, LIGHT_BLUE),
    ('Improve\nSelf-Assessment\n(LCAE)', '學姐已證', TEAL, LIGHT_TEAL),
    ('More Accurate\nToken Need\nPrediction', '本研究驗證', RED_ACCENT, LIGHT_RED),
    ('More Efficient\nToken\nAllocation', '本研究驗證', RED_ACCENT, LIGHT_RED),
]

box_w = Emu(2400000)
box_h = Emu(1645920)
y = Emu(2286000)

for i, (text, proof, color, fill) in enumerate(chain_items):
    x = Emu(457200 + (2400000 + 457200) * i)
    add_rounded_rect(slide, x, y, box_w, box_h, fill_color=fill, line_color=color, line_width=Pt(2))
    add_textbox(slide, Emu(x.emu + 137160), Emu(y.emu + 228600), Emu(2100000), Emu(822960),
                text, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Proof label below
    add_textbox(slide, Emu(x.emu + 137160), Emu(y.emu + 1180000), Emu(2100000), Emu(365760),
                proof, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 3:
        arrow_x = Emu(x.emu + box_w.emu + 91440)
        add_arrow(slide, arrow_x, Emu(y.emu + 548640), Emu(274320), Emu(548640))

# Story text
add_shape(slide, Emu(457200), Emu(4297680), Emu(11277600), Emu(1828800), fill_color=LIGHT_GRAY)
add_textbox(slide, Emu(548640), Emu(4343400), Emu(10972800), Emu(365760),
            'One-sentence story', font_size=16, bold=True, color=DARK_BLUE)
add_textbox(slide, Emu(548640), Emu(4754880), Emu(10972800), Emu(1180000),
            "學姐 proved IRT difficulty signals improve LLM self-assessment (LCAE) and mentioned reliability-inference cost association. "
            "This research fills that gap: we measure the self-assessment - token efficiency relationship, validate the "
            "calibration -> token savings causal chain, and propose a self-assessment-based token allocation strategy. "
            "If it holds, this provides a training-free token optimization method - no extra training needed, just let the model know itself better.",
            font_size=14, color=DARK_TEXT)

add_footer(slide, 10)
add_notes(slide, "這是整個論文的故事。四個方塊從左到右是一條因果鏈。前兩步學姐已經證明了：IRT 難度信號能改善自我評估。後兩步是本研究要驗證的：更準的自我評估能不能預測 token 需求、能不能做到更高效的 token 分配。底下是一句話版本的故事：如果成立，這提供了一個 training-free 的 token 優化方法，不需要額外訓練，只需要讓模型更了解自己。這是跟 ROI-Reasoning 等競爭者最大的差別——他們需要訓練 predictor，我們用模型本身的自我評估。")


# ============================================================
# SLIDE 11: 四個貢獻
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Contributions (4)', 'What this paper adds to the field')
add_bottom_bar(slide)

contribs = [
    ('C1', 'New Finding', 'Reveal the relationship between\nself-assessment accuracy (LCAE)\nand token requirements', 'Gap 1', RED_ACCENT, LIGHT_RED),
    ('C2', 'New Method', 'Use IRT difficulty signals\nto predict token requirements', 'Gap 2', AMBER, LIGHT_AMBER),
    ('C3', 'New Validation', 'Validate the causal chain:\nbetter calibration -> less tokens\nfor same quality', 'Gap 4', RED_ACCENT, LIGHT_RED),
    ('C4', 'New Framework', 'Propose self-assessment-based\ntoken allocation strategy\n(training-free)', 'Gap 1+2+4', TEAL, LIGHT_TEAL),
]

card_w = Emu(2700000)
card_h = Emu(3657600)
card_y = Emu(1463040)

for i, (label, ctype, desc, gap, border, fill) in enumerate(contribs):
    x = Emu(457200 + (2700000 + 228600) * i)
    add_left_border_card(slide, x, card_y, card_w, card_h, border_color=border, fill_color=fill)
    # Label circle
    add_shape(slide, Emu(x.emu + 137160), Emu(card_y.emu + 91440), Emu(548640), Emu(548640),
              fill_color=border, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # Type
    add_textbox(slide, Emu(x.emu + 822960), Emu(card_y.emu + 137160), Emu(1700000), Emu(365760),
                ctype, font_size=14, bold=True, color=border)
    # Description
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 822960), Emu(2400000), Emu(1645920),
                desc, font_size=14, color=DARK_TEXT)
    # Gap label
    add_shape(slide, Emu(x.emu + 137160), Emu(card_y.emu + 2743200), Emu(1645920), Emu(457200),
              fill_color=border)
    add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 2743200 + 45720), Emu(1400000), Emu(365760),
                f'Addresses {gap}', font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 11)
add_notes(slide, "四個貢獻。C1 是新發現：揭示 LCAE 和 token 需求的關係。C2 是新方法：用 IRT 難度信號預測 token 需求。C3 是新驗證：驗證校準好到 token 省的因果鏈。C4 是新框架：提出基於自我評估的 token 分配策略，而且是 training-free 的。四個貢獻涵蓋了發現、方法、驗證、框架四個層面，對應前面講的四個 gap。")


# ============================================================
# SLIDE 12: 競爭者比較
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Nearest Competitors', 'Why nobody has done this before')
add_bottom_bar(slide)

# Table
headers = ['Competitor', 'Overlap', 'Key Difference']
rows = [
    ['ROI-Reasoning\n(2026-01)', 'Predicts reasoning cost\n+ budget allocation', 'No IRT, no self-assessment,\nuses trained predictor'],
    ['TRIAGE\n(2026-05)', 'Metacognitive control\n+ token budget', 'Evaluates allocation quality,\nnot calibration x token'],
    ['學姐 LCAE\n(IEEE IRI 2026)', 'IRT + calibration\n+ mentions cost', 'Does not go deep\ninto token efficiency'],
    ['UAB\n(2026-05)', 'Uncertainty ->\nbudget allocation', 'Uses log-prob, not\nself-assessment'],
    ['"LLM Already Knows"\n(2025-09)', 'Pre-generation\ndifficulty prediction', 'Uses hidden states, not\nverbalized confidence'],
]

add_table_simple(slide, Emu(457200), Emu(1463040), Emu(11277600), Emu(4114800),
                 headers, rows,
                 col_widths=[Emu(2743200), Emu(4114800), Emu(4419600)])

# Key takeaway
add_shape(slide, Emu(457200), Emu(5669280), Emu(11277600), Emu(548640), fill_color=DARK_BLUE)
add_textbox(slide, Emu(548640), Emu(5715240), Emu(10972800), Emu(365760),
            "No paper simultaneously does ALL FOUR: IRT framework + self-assessment measurement + token prediction + causal validation",
            font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 12)
add_notes(slide, "五個最接近的競爭者。ROI-Reasoning 做了 meta-cognitive 預測加 budget 分配，但不用 IRT 也不用自我評估，是自訓 predictor。TRIAGE 測 metacognitive budget 但不研究校準和 token 的關係。學姐有 IRT 加校準但沒深入 token。UAB 用 log-prob 不是自我評估。LLM Already Knows 用 hidden states 不是 verbalized confidence。重點是沒有任何一篇同時做到四件事：IRT 框架、自我評估測量、token 預測、因果驗證。")


# ============================================================
# SLIDE 13: 實驗設計總覽
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Experiment Design Overview', 'Three phases, building on 學姐 assets')
add_bottom_bar(slide)

# Phase flow
phases = [
    ('Phase 1', 'Correlation Analysis', 'Use 學姐 20 models data\n+ measure token usage\n-> LCAE x Token efficiency\ncorrelation', DARK_BLUE, LIGHT_BLUE),
    ('Phase 2', 'Causal Chain Validation', 'With/without IDS ( 學姐 proved\nimproves LCAE) -> measure\ntoken change -> causal effect\nof calibration on tokens', TEAL, LIGHT_TEAL),
    ('Phase 3', 'Token Allocation Strategy', 'Design allocation based on\nself-assessment -> compare\nwith Uniform/Oracle/External\n-> measure efficiency', RED_ACCENT, LIGHT_RED),
]

card_w = Emu(3403520)
card_h = Emu(3200400)
card_y = Emu(1463040)

for i, (phase, title, desc, color, fill) in enumerate(phases):
    x = Emu(457200 + (3403520 + 457200) * i)
    add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=fill, line_color=color, line_width=Pt(2))
    # Phase label
    add_shape(slide, Emu(x.emu + 137160), Emu(card_y.emu + 91440), Emu(1180000), Emu(457200), fill_color=color)
    add_textbox(slide, Emu(x.emu + 228600), Emu(card_y.emu + 137160), Emu(1000000), Emu(320040),
                phase, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 640080), Emu(3100000), Emu(457200),
                title, font_size=18, bold=True, color=color)
    # Description
    add_textbox(slide, Emu(x.emu + 137160), Emu(card_y.emu + 1180000), Emu(3100000), Emu(1828800),
                desc, font_size=14, color=DARK_TEXT)
    # Arrow
    if i < 2:
        arrow_x = Emu(x.emu + card_w.emu + 91440)
        add_arrow(slide, arrow_x, Emu(card_y.emu + 1371600), Emu(274320), Emu(548640))

# Bottom: shared resources
add_shape(slide, Emu(457200), Emu(4983480), Emu(11277600), Emu(1097280), fill_color=LIGHT_GRAY)
add_textbox(slide, Emu(548640), Emu(5029200), Emu(10972800), Emu(365760),
            'Shared Resources (from 學姐)', font_size=16, bold=True, color=DARK_BLUE)
add_textbox(slide, Emu(548640), Emu(5392680), Emu(10972800), Emu(548640),
            "20 models IRT ability estimates (theta)  |  Item difficulty parameters (b)  |  LCAE calculation method  "
            "|  IDS/QOQ/DPR/Combined self-assessment scenarios  |  Benchmark data  |  Prompt templates",
            font_size=13, color=DARK_TEXT)

add_footer(slide, 13)
add_notes(slide, "實驗分三個階段。Phase 1 用學姐的 20 個模型資料，在同一 benchmark 上測量每個模型的 token 使用量，分析 LCAE 和 token 效率的相關性。Phase 2 做因果鏈驗證：學姐已經證明 IDS 能改善 LCAE，我在有無 IDS 的情況下測量 token 變化，如果 IDS 同時改善了 token 效率而且改善幅度跟 LCAE 改善幅度正相關，就支持因果鏈。Phase 3 設計分配策略：讓模型自我評估後根據信心分配 token，跟 Uniform、Oracle、外部模型等對照組比較。底下列的是從學姐那裡可以直接重用的資源。")


# ============================================================
# SLIDE 14: Phase 1 詳細
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Phase 1: LCAE x Token Correlation', 'Does self-assessment accuracy predict token efficiency?')
add_bottom_bar(slide)

# Left: Variables
add_textbox(slide, Emu(457200), Emu(1463040), Emu(5486400), Emu(457200),
            'Measured Variables', font_size=18, bold=True, color=DARK_BLUE)

var_data = [
    ['LCAE', 'from 學姐', 'Self-assessment vs IRT consistency'],
    ['Ability (theta)', 'from 學姐', 'IRT estimated model capability'],
    ['Difficulty (b)', 'from 學姐', 'IRT estimated item difficulty'],
    ['Token usage (T)', 'NEW', 'Actual reasoning tokens per question'],
    ['Accuracy', 'from 學姐', 'Answer correctness'],
    ['Token efficiency', 'NEW', 'Accuracy / Token usage'],
]

add_table_simple(slide, Emu(457200), Emu(2011680), Emu(5486400), Emu(2926080),
                 ['Variable', 'Source', 'Description'], var_data,
                 col_widths=[Emu(1645920), Emu(1371600), Emu(2468880)])

# Right: Steps
add_textbox(slide, Emu(6217920), Emu(1463040), Emu(5486400), Emu(457200),
            'Steps', font_size=18, bold=True, color=DARK_BLUE)

steps = [
    '1. Run 20 models on 學姐 benchmark,\n   measure reasoning tokens per question',
    '2. Compute LCAE vs avg token efficiency\ncorrelation (Pearson/Spearman)',
    '3. Stratified analysis by difficulty:\n   easy / medium / hard',
    '4. Analyze: do calibrated models show\n   adaptive allocation (easy=less,\n   hard=more)?',
]

y = Emu(2011680)
for step in steps:
    add_left_border_card(slide, Emu(6217920), Emu(y), Emu(5486400), Emu(731520), border_color=TEAL)
    add_textbox(slide, Emu(6400800), Emu(y + 91440), Emu(5200000), Emu(548640),
                step, font_size=13, color=DARK_TEXT)
    y += 822960

# Expected results
add_shape(slide, Emu(457200), Emu(5212320), Emu(11277600), Emu(914400), fill_color=LIGHT_TEAL)
add_shape(slide, Emu(457200), Emu(5212320), Emu(54864), Emu(914400), fill_color=TEAL)
add_textbox(slide, Emu(548640), Emu(5257680), Emu(10972800), Emu(365760),
            'Expected Results', font_size=14, bold=True, color=TEAL)
add_textbox(slide, Emu(548640), Emu(5623920), Emu(10972800), Emu(457200),
            "H1 holds: LCAE positively correlates with token efficiency, especially on hard questions. "
            "H1 null: self-assessment and token efficiency are unrelated (also an important finding).",
            font_size=13, color=DARK_TEXT)

add_footer(slide, 14)
add_notes(slide, "Phase 1 的詳細設計。左邊是測量變量：LCAE、能力、難度都從學姐資料來，token 使用量是新增的，token efficiency 是新定義的——accuracy 除以 token usage。右邊是四個步驟：跑 20 個模型測 token、算相關性、按難度分層分析、看校準好的模型是否表現出適應性分配。預期結果有兩種：H1 成立就是 LCAE 和 token 效率正相關，尤其在困難題上更顯著；如果不成立，自我評估和 token 效率無關本身也是重要發現。控制變因包括固定 temperature、max_tokens 足夠大不截斷、同一 prompt template。")


# ============================================================
# SLIDE 15: Phase 2 詳細
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Phase 2: Causal Chain Validation', 'Does improving calibration CAUSE token efficiency gains?')
add_bottom_bar(slide)

# Experiment groups table
add_textbox(slide, Emu(457200), Emu(1463040), Emu(5486400), Emu(457200),
            'Experiment Groups (學姐 proved LCAE differences)', font_size=16, bold=True, color=DARK_BLUE)

group_data = [
    ['Control', 'None (QOQ)', 'Worst', 'Token + Accuracy'],
    ['Treatment 1', 'IDS', 'Better', 'Token + Accuracy'],
    ['Treatment 2', 'DPR', 'Medium', 'Token + Accuracy'],
    ['Treatment 3', 'Combined', 'Best', 'Token + Accuracy'],
]

add_table_simple(slide, Emu(457200), Emu(2011680), Emu(5486400), Emu(2194560),
                 ['Group', 'Signal', 'Expected LCAE', 'Measure'], group_data,
                 col_widths=[Emu(1371600), Emu(1371600), Emu(1371600), Emu(1371600)])

# Causal logic
add_textbox(slide, Emu(6217920), Emu(1463040), Emu(5486400), Emu(457200),
            'Causal Logic', font_size=16, bold=True, color=DARK_BLUE)

logic_steps = [
    'IDS improves LCAE ( 學姐 proved)',
    'If IDS also improves token efficiency...',
    'And improvement magnitude correlates\n  with LCAE improvement...',
    'Then: calibration -> token savings\n  causal chain is supported',
]

y = Emu(2011680)
for i, step in enumerate(logic_steps):
    add_left_border_card(slide, Emu(6217920), Emu(y), Emu(5486400), Emu(548640),
                         border_color=TEAL if i < 3 else RED_ACCENT)
    add_textbox(slide, Emu(6400800), Emu(y + 45720), Emu(5200000), Emu(457200),
                step, font_size=13, color=DARK_TEXT)
    y += 640080

# Further analysis
add_textbox(slide, Emu(457200), Emu(4389120), Emu(11277600), Emu(365760),
            'Further Analysis', font_size=16, bold=True, color=DARK_BLUE)

analysis_items = [
    'Difficulty stratification: In which difficulty range does IDS have the largest token effect?',
    'Allocation quality: Does IDS group show adaptive allocation (easy=less, hard=more)?',
    'Overthinking analysis: Does IDS reduce redundant reasoning on easy questions?',
    'Control: Free budget (natural) vs Fixed budget (constrained) vs Adaptive budget (self-assessed)',
]

y = Emu(4754880)
for item in analysis_items:
    add_shape(slide, Emu(457200), Emu(y + 91440), Emu(91440), Emu(91440), fill_color=TEAL, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Emu(640080), Emu(y), Emu(11100000), Emu(457200),
                item, font_size=13, color=DARK_TEXT)
    y += 457200

add_footer(slide, 15)
add_notes(slide, "Phase 2 是因果鏈驗證。利用學姐已經證明的結果：IDS 能改善 LCAE，四種情境的 LCAE 好壞順序是已知的。如果 IDS 組同時也改善了 token 效率，而且改善幅度跟 LCAE 改善幅度正相關，就支持因果鏈。進一步分析包括按難度分層看 IDS 在哪個範圍影響最大、IDS 組是否表現出適應性分配、是否減少了簡單題的 overthinking。控制實驗有三種 budget 條件：free budget 讓模型自然決定、fixed budget 限制總量看分配品質、adaptive budget 根據自我評估分配。")


# ============================================================
# SLIDE 16: Phase 3 詳細
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Phase 3: Token Allocation Strategy', 'Can self-assessment-based allocation outperform baselines?')
add_bottom_bar(slide)

# Strategy flow
add_textbox(slide, Emu(457200), Emu(1463040), Emu(11277600), Emu(365760),
            'Allocation Pipeline', font_size=16, bold=True, color=DARK_BLUE)

flow_steps = [
    ('Model receives\nquestion', DARK_BLUE, LIGHT_BLUE),
    ('Self-assessment\n(confidence/difficulty)', TEAL, LIGHT_TEAL),
    ('Allocate token\nbudget based on\nself-assessment', AMBER, LIGHT_AMBER),
    ('Generate answer\nwithin budget', RED_ACCENT, LIGHT_RED),
    ('Evaluate accuracy\n+ token efficiency', RGBColor(0x6A, 0x4C, 0x93), RGBColor(0xF0, 0xEC, 0xF8)),
]

box_w = Emu(2100000)
box_h = Emu(1180000)
y = Emu(2011680)

for i, (text, color, fill) in enumerate(flow_steps):
    x = Emu(457200 + (2100000 + 228600) * i)
    add_rounded_rect(slide, x, y, box_w, box_h, fill_color=fill, line_color=color, line_width=Pt(1.5))
    add_textbox(slide, Emu(x.emu + 91440), Emu(y.emu + 228600), Emu(1900000), Emu(731520),
                text, font_size=13, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    if i < 4:
        arrow_x = Emu(x.emu + box_w.emu + 45720)
        add_arrow(slide, arrow_x, Emu(y.emu + 411480), Emu(137160), Emu(365760))

# Comparison table
add_textbox(slide, Emu(457200), Emu(3429360), Emu(11277600), Emu(365760),
            'Comparison Strategies', font_size=16, bold=True, color=DARK_BLUE)

strat_data = [
    ['Uniform', 'Same budget for all', 'Baseline'],
    ['Oracle', 'IRT difficulty (post-hoc)', 'Upper bound'],
    ['Self-assessment (verbalized)', 'Verbalized confidence', 'Tests H4'],
    ['Self-assessment (log-prob)', 'Log-prob confidence', 'Tests H4'],
    ['IDS-based', 'IRT signal + self-assessment', 'Combines 學姐 framework'],
    ['External estimator', 'External model predicts', 'ROI-Reasoning style'],
]

add_table_simple(slide, Emu(457200), Emu(3885600), Emu(11277600), Emu(2194560),
                 ['Strategy', 'Token allocation based on', 'Role'], strat_data,
                 col_widths=[Emu(3403520), Emu(4114800), Emu(3759240)])

# Metrics
add_textbox(slide, Emu(457200), Emu(6172200), Emu(5486400), Emu(365760),
            'Metrics: Token efficiency | Budget utilization | Allocation quality (Spearman rho vs Oracle) | Regret',
            font_size=12, color=MEDIUM_GRAY)

add_footer(slide, 16)
add_notes(slide, "Phase 3 是分配策略。Pipeline 是：模型收到題目、做自我評估、根據信心分配 token 預算、在預算內生成回答、評估品質和效率。對照組有六個：Uniform 是所有題目相同預算做基線，Oracle 是用 IRT 難度事後分配做上限，verbalized 和 log-prob 兩種自我評估測 H4，IDS-based 結合學姐框架，External estimator 對照 ROI-Reasoning 風格。評估指標包括 token efficiency、budget utilization、allocation quality 跟 Oracle 的相關性、以及 regret。")


# ============================================================
# SLIDE 17: 與學姐論文的銜接
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Integration with 學姐 Framework', 'What we reuse vs what we add')
add_bottom_bar(slide)

# Left: Reuse
add_left_border_card(slide, Emu(457200), Emu(1463040), Emu(5303520), Emu(4297680),
                     border_color=GREEN_ACCENT, fill_color=LIGHT_GREEN)
add_textbox(slide, Emu(640080), Emu(1554480), Emu(4937760), Emu(457200),
            'Reused from 學姐', font_size=18, bold=True, color=GREEN_ACCENT)

reuse_items = [
    '20 models IRT ability estimates (theta)',
    'Item difficulty parameters (b)',
    'LCAE metric calculation method',
    'IDS/QOQ/DPR/Combined scenarios',
    'Self-assessment prompt templates',
    'Benchmark data',
]

y = Emu(2194560)
for item in reuse_items:
    add_textbox(slide, Emu(822960), Emu(y), Emu(4754880), Emu(457200),
                f'  {item}', font_size=14, color=DARK_TEXT)
    y += 548640

# Right: New
add_left_border_card(slide, Emu(6217920), Emu(1463040), Emu(5303520), Emu(4297680),
                     border_color=RED_ACCENT, fill_color=LIGHT_RED)
add_textbox(slide, Emu(6400800), Emu(1554480), Emu(4937760), Emu(457200),
            'New in This Research', font_size=18, bold=True, color=RED_ACCENT)

new_items = [
    'Token usage precise measurement',
    'Token efficiency metric definition',
    'Causal chain validation experiment',
    'Token allocation strategy design',
    'Verbalized vs log-prob comparison',
    'Overthinking analysis',
]

y = Emu(2194560)
for item in new_items:
    add_textbox(slide, Emu(6583680), Emu(y), Emu(4754880), Emu(457200),
                f'  {item}', font_size=14, color=DARK_TEXT)
    y += 548640

# Narrative connection
add_shape(slide, Emu(457200), Emu(5852160), Emu(11277600), Emu(457200), fill_color=DARK_BLUE)
add_textbox(slide, Emu(548640), Emu(5852160), Emu(10972800), Emu(457200),
            "學姐 story: IRT -> LCAE -> IDS improves calibration -> mentions cost association   ||   "
            "Ryan story: LCAE -> token prediction -> token allocation optimization -> validate causal chain",
            font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 17)
add_notes(slide, "左邊是從學姐那裡直接重用的：20 個模型的 IRT 估計、題目難度參數、LCAE 計算方法、四種自我評估情境、prompt templates、benchmark 資料。這些都不需要重新做。右邊是本研究新增的：token 精確測量、token efficiency 定義、因果實驗、分配策略、verbalized vs log-prob 比較、overthinking 分析。底部的敘事銜接：學姐的故事是 IRT 到 LCAE 到 IDS 改善校準到提到 cost 關聯，我的故事是 LCAE 到 token 預測到 token 分配優化到驗證因果鏈。把學姐提到的 cost 關聯變成深入驗證的因果鏈。")


# ============================================================
# SLIDE 18: 風險評估
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Risk Assessment', 'Risks and mitigation strategies')
add_bottom_bar(slide)

risks = [
    ('High', 'Verbalized confidence may not work for token allocation',
     'Kumaran (2026-06) found it tracks commitment, not correctness',
     'Test both verbalized and log-prob; frame as research question, not assumption; even if verbalized fails, IDS may still work',
     RED_ACCENT, LIGHT_RED),
    ('High', 'Causal chain may not hold',
     'Better calibration might not cause token savings',
     'Null result is itself a finding; distinguish allocation quality from total volume; report honestly',
     RED_ACCENT, LIGHT_RED),
    ('Medium', 'Scope may be too narrow',
     'Only one intersection: self-assessment x token',
     'Broaden to metacognition x compute allocation; 4 contributions cover finding, method, validation, framework',
     AMBER, LIGHT_AMBER),
    ('Medium', ' 學姐 data may not have token records',
     'Need to re-run experiments',
     'Confirm data format first; if needed, run representative subset (5-8 models)',
     AMBER, LIGHT_AMBER),
    ('Low', 'ROI-Reasoning too close',
     'Already does meta-cognitive prediction + budget',
     'Differences sufficient: no IRT, no calibration, no causal analysis, uses trained predictor not self-assessment',
     TEAL, LIGHT_TEAL),
]

y = Emu(1463040)
for level, risk, detail, mitigation, color, fill in risks:
    add_left_border_card(slide, Emu(457200), Emu(y), Emu(11277600), Emu(914400),
                         border_color=color, fill_color=fill)
    # Level badge
    add_shape(slide, Emu(548640), Emu(y + 91440), Emu(822960), Emu(365760), fill_color=color)
    add_textbox(slide, Emu(548640), Emu(y + 91440 + 45720), Emu(822960), Emu(274320),
                level, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Risk
    add_textbox(slide, Emu(1554480), Emu(y + 91440), Emu(4800000), Emu(320040),
                risk, font_size=14, bold=True, color=DARK_TEXT)
    # Detail
    add_textbox(slide, Emu(1554480), Emu(y + 411480), Emu(4800000), Emu(365760),
                detail, font_size=12, color=MEDIUM_GRAY)
    # Mitigation
    add_textbox(slide, Emu(6494040), Emu(y + 91440), Emu(5100000), Emu(731520),
                f'Mitigation: {mitigation}', font_size=12, color=DARK_TEXT)
    y += 1000000

add_footer(slide, 18)
add_notes(slide, "五個風險。兩個高風險：第一是 verbalized confidence 可能不適合做 token 分配信號，Kumaran 發現它追蹤 commitment 不是 correctness。緩解是同時測 verbalized 和 log-prob，把它當研究問題而非假設，而且 IDS 不依賴 verbalized confidence。第二是因果鏈可能不成立，校準好不代表 token 省。緩解是不成立本身也是發現，而且要區分分配品質和總量。兩個中風險：範圍太窄和學姐資料可能沒有 token 記錄。一個低風險是 ROI-Reasoning 太接近，但差異足夠大。")


# ============================================================
# SLIDE 19: 時程規劃
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Timeline (6 weeks to IEEE Big Data deadline)', 'Aug 12, 2026 submission target')
add_bottom_bar(slide)

# Timeline table
timeline_data = [
    ['W1', 'Jul 1-7', 'Confirm 學姐 data format, get 20 model list, check token records', 'Data checklist'],
    ['W2', 'Jul 8-14', 'Phase 1 data collection (run models + token measurement)', 'Phase 1 raw data'],
    ['W3', 'Jul 15-21', 'Phase 1 analysis + Phase 2 experiment execution', 'Phase 1 results + Phase 2 data'],
    ['W4', 'Jul 22-28', 'Phase 2 analysis + Phase 3 strategy design + experiment', 'Phase 2 results + Phase 3 data'],
    ['W5', 'Jul 29-Aug 4', 'Phase 3 analysis + paper draft writing', 'Full results + draft'],
    ['W6', 'Aug 5-11', 'Paper revision, supplementary experiments, finalization', 'Final paper'],
]

add_table_simple(slide, Emu(457200), Emu(1463040), Emu(11277600), Emu(2743200),
                 ['Week', 'Date', 'Task', 'Deliverable'], timeline_data,
                 col_widths=[Emu(822960), Emu(1645920), Emu(5486400), Emu(3296280)])

# Degradation plan
add_textbox(slide, Emu(457200), Emu(4297680), Emu(11277600), Emu(365760),
            'Degradation Plan (if time runs short)', font_size=16, bold=True, color=DARK_BLUE)

levels = [
    ('Level 0', 'Full: Phase 1+2+3', 'C1+C2+C3+C4', GREEN_ACCENT),
    ('Level 1', 'Phase 1+2, Phase 3 simplified (design only)', 'C1+C2+C3', TEAL),
    ('Level 2', 'Phase 1+2 only', 'C1+C3 (correlation + causation)', AMBER),
    ('Level 3', 'Phase 1 only', 'C1 (correlation finding)', RED_ACCENT),
]

y = Emu(4754880)
for level, desc, contribs, color in levels:
    add_shape(slide, Emu(457200), Emu(y + 45720), Emu(1645920), Emu(320040), fill_color=color)
    add_textbox(slide, Emu(457200), Emu(y + 45720 + 45720), Emu(1645920), Emu(228600),
                level, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Emu(2286000), Emu(y + 45720), Emu(5486400), Emu(320040),
                desc, font_size=13, color=DARK_TEXT)
    add_textbox(slide, Emu(8229600), Emu(y + 45720), Emu(3403520), Emu(320040),
                contribs, font_size=12, bold=True, color=color)
    y += 411480

# Note
add_textbox(slide, Emu(457200), Emu(6400800), Emu(5486400), Emu(228600),
            'Even Level 3 is publishable: first discovery of LCAE-token relationship', font_size=12, color=MEDIUM_GRAY)

add_footer(slide, 19)
add_notes(slide, "六週時程。第一週確認學姐資料格式和模型清單。第二週跑 Phase 1 收集 token 資料。第三週 Phase 1 分析加 Phase 2 實驗。第四週 Phase 2 分析加 Phase 3。第五週分析和寫初稿。第六週修改定稿。如果時間不夠有四級降級方案。Level 0 是全部做完，Level 1 是 Phase 3 簡化，Level 2 只做 Phase 1 加 2，Level 3 只做 Phase 1。即使 Level 3 也還是有貢獻的——首次揭示 LCAE 和 token 的關係，可以寫 short paper 或 workshop paper。")


# ============================================================
# SLIDE 20: 給老師的問題
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_title_bar(slide, 'Questions for Advisor', '8 questions for discussion')
add_bottom_bar(slide)

questions = [
    ('Q1', 'Direction', 'Is this story - self-assessment -> token prediction -> allocation optimization - viable?'),
    ('Q2', 'Data access', 'Does 學姐 data include token records? Can I access the 20-model data and benchmark?'),
    ('Q3', 'Venue', 'Invest IEEE Big Data 2026 (Aug deadline) or do deeper work for another venue?'),
    ('Q4', 'Priority', 'If time is short, which Phase to cut? P1 (correlation) / P2 (causal) / P3 (strategy)?'),
    ('Q5', 'Benchmark', 'What benchmark did 學姐 use? Need to expand to multi-domain (MMLU + GSM8K)?'),
    ('Q6', 'Verbalized conf', 'Is Kumaran finding a risk or a research opportunity? Frame as question or assume log-prob?'),
    ('Q7', 'nhri project', 'Combine with memory conflict research? Or focus on one direction?'),
    ('Q8', 'Models', 'Re-run all 20 models or select representative subset (5-8, covering frontier/mid/small)?'),
]

y = Emu(1463040)
for label, topic, question in questions:
    # Number badge
    add_shape(slide, Emu(457200), Emu(y), Emu(548640), Emu(548640),
              fill_color=DARK_BLUE, shape_type=MSO_SHAPE.OVAL)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    # Topic
    add_textbox(slide, Emu(1180000), Emu(y + 22860), Emu(2400000), Emu(320040),
                topic, font_size=13, bold=True, color=TEAL)
    # Question
    add_textbox(slide, Emu(3800000), Emu(y + 22860), Emu(7800000), Emu(457200),
                question, font_size=14, color=DARK_TEXT)
    y += 640080

add_footer(slide, 20)
add_notes(slide, "八個問題想跟老師討論。最重要的幾個：Q1 方向是否可行；Q2 學姐資料有沒有 token 記錄、能不能拿到；Q3 投 IEEE Big Data 還是做更深；Q4 時間不夠砍哪個 Phase。Q5 是 benchmark 選擇，Q6 是 verbalized confidence 的風險還是機會，Q7 是要不要跟 nhri 記憶衝突研究結合，Q8 是 20 模型全跑還是選子集。")


# ============================================================
# SLIDE 21: Summary / Thank You
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)

# Top and bottom accent
add_shape(slide, Emu(0), Emu(0), SLIDE_W, Emu(54864), fill_color=TEAL)
add_shape(slide, Emu(0), Emu(6537960), SLIDE_W, Emu(320040), fill_color=DARK_BLUE)

# Center content
add_textbox(slide, Emu(1371600), Emu(1828800), Emu(9448495), Emu(640080),
            'Summary', font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

summary_points = [
    " 學姐 proved IRT difficulty signals improve LLM self-assessment (LCAE)",
    "She mentioned reliability-cost association but did not explore it",
    "16 systematic searches confirm: nobody connects calibration to token efficiency",
    "This research validates the causal chain: calibration -> token prediction -> allocation",
    "4 contributions: new finding, method, validation, framework",
    "Training-free optimization: no extra training, just better self-awareness",
]

y = Emu(2743200)
for point in summary_points:
    add_shape(slide, Emu(2743200), Emu(y + 45720), Emu(91440), Emu(91440), fill_color=TEAL, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Emu(2926080), Emu(y), Emu(8229600), Emu(457200),
                point, font_size=15, color=DARK_TEXT)
    y += 457200

# Thank you
add_textbox(slide, Emu(3657600), Emu(5486400), Emu(4876495), Emu(548640),
            'Thank you - Questions & Discussion', font_size=20, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

add_footer(slide, 21)
add_notes(slide, "總結六個重點。學姐證明了 IRT 改善校準、提到了 cost 關聯但沒深入。16 組搜尋確認沒有人把校準和 token 效率連起來。本研究驗證因果鏈、提出分配策略、四個貢獻涵蓋發現方法驗證框架。最大賣點是 training-free：不需要額外訓練，只需要讓模型更了解自己。謝謝老師，歡迎討論。")


# ============================================================
# Save
# ============================================================
output_path = '/Users/ryan/Projects/llm-calibration-token-efficiency/slides.pptx'
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")