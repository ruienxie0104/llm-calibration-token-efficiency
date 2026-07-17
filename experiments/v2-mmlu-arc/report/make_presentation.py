#!/usr/bin/env python3
"""Generate weekly report presentation for NHRI advisor meeting."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

FIG = "experiments/v2-mmlu-arc/report/figures"
FIG_V1 = "experiments/v1-pilot/report/figures"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x4a, 0x6f, 0xd8)
ACCENT = RGBColor(0xe8, 0x7a, 0x3c)
GREEN = RGBColor(0x4a, 0xa8, 0x5c)
RED = RGBColor(0xc0, 0x39, 0x2b)
GRAY = RGBColor(0x6b, 0x6b, 0x80)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG = RGBColor(0xf5, 0xf5, 0xfa)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xa0, 0xa0, 0xb0)
    p2.alignment = PP_ALIGN.CENTER
    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(0.5))
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = "2026-07-17"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xd0, 0xd0, 0xe0)
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, image=None, image_caption=None, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    add_title_bar(slide, title)
    # Content
    left = Inches(0.6)
    top = Inches(1.3)
    width = Inches(12.13)
    height = Inches(5.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if level == 0 else Pt(14)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.level = level
        p.space_after = Pt(6)
    # Image
    if image and os.path.exists(image):
        pic = slide.shapes.add_picture(image, Inches(7), Inches(1.5), Inches(6), Inches(5.5))
        if image_caption:
            capBox = slide.shapes.add_textbox(Inches(7), Inches(7), Inches(6), Inches(0.4))
            capP = capBox.text_frame.paragraphs[0]
            capP.text = image_caption
            capP.font.size = Pt(11)
            capP.font.color.rgb = GRAY
            capP.alignment = PP_ALIGN.CENTER
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_image_slide(prs, title, image_path, caption="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    if os.path.exists(image_path):
        # Center the image
        pic = slide.shapes.add_picture(image_path, Inches(2), Inches(1.4), Inches(9.33), Inches(5.5))
        if caption:
            capBox = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.33), Inches(0.5))
            capP = capBox.text_frame.paragraphs[0]
            capP.text = caption
            capP.font.size = Pt(12)
            capP.font.color.rgb = GRAY
            capP.alignment = PP_ALIGN.CENTER
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

def add_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.45 * rows_count))
    table = table_shape.table
    # Headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    # Rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def add_two_image_slide(prs, title, img1, img2, caption1="", caption2=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    if os.path.exists(img1):
        slide.shapes.add_picture(img1, Inches(0.3), Inches(1.4), Inches(6.3), Inches(5))
    if os.path.exists(img2):
        slide.shapes.add_picture(img2, Inches(6.8), Inches(1.4), Inches(6.3), Inches(5))
    if caption1:
        c1 = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(6.3), Inches(0.4))
        p1 = c1.text_frame.paragraphs[0]
        p1.text = caption1
        p1.font.size = Pt(11)
        p1.font.color.rgb = GRAY
        p1.alignment = PP_ALIGN.CENTER
    if caption2:
        c2 = slide.shapes.add_textbox(Inches(6.8), Inches(6.5), Inches(6.3), Inches(0.4))
        p2 = c2.text_frame.paragraphs[0]
        p2.text = caption2
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

def add_title_bar(slide, title):
    """Add a colored title bar at the top."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK
    # Underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(12.83), Inches(1.1))
    line.line.color.rgb = BLUE
    line.line.width = Pt(2)

# ============================================================
# Build slides
# ============================================================

# Slide 1: Title
add_title_slide(prs,
    "PM × LLM 推理路徑分析：Pilot 階段分享",
    "Calibration-Aware Token Efficiency — V1 & V2 實驗結果")

# Slide 2: Outline
add_content_slide(prs, "大綱", [
    ("1. 上次報告回顧與回饋", 0),
    ("   • 7/4 提案：校準 → Token 分配品質的因果鏈", 1),
    ("   • 老師回饋與方向調整", 1),
    ("2. V1 Pilot 實驗", 0),
    ("   • GSM8K, 20 題, 5 模型", 1),
    ("   • 發現：PM 能區分推理風格，但題目太簡單", 1),
    ("3. V2 實驗設計與執行", 0),
    ("   • MMLU STEM + ARC Challenge, 100 題, 4 模型", 1),
    ("   • 新增：信心自評、校準指標", 1),
    ("4. V2 關鍵發現", 0),
    ("   • 信心差距是關鍵校準指標", 1),
    ("   • Petri net 一致性檢查的局限", 1),
    ("   • 熵分析 + 步驟類型散度（新方法）", 1),
    ("5. 下一步", 0),
])

# Slide 3: Section - 上次報告回顧
add_section_slide(prs, "1. 上次報告回顧")

# Slide 4: 上次提案核心
add_content_slide(prs, "上次提案：校準 → Token 分配品質", [
    ("研究核心問題", 0),
    ("  學姐證明 IRT 難度信號 (IDS) 能改善 LLM 信心校準", 1),
    ("  但沒有驗證：校準改善後，token 分配品質是否更好？", 1),
    ("", 0),
    ("三個研究問題", 0),
    ("  RQ1: LCAE 能否預測 excess token usage？", 1),
    ("  RQ2: IDS 能否改善 token allocation quality？", 1),
    ("  RQ3: 校準改善 → allocation quality 提升的因果鏈成立？", 1),
    ("", 0),
    ("賣點：Training-free — 不需額外訓練，只靠模型自我評估", 0),
    ("", 0),
    ("文獻調查：16 組 arXiv 搜尋，五個集群彼此不交集", 0),
    ("  沒有人把 自我評估 + IRT + Token 分配 連起來", 1),
], notes="這是 7/4 的提案核心。學姐的框架到 LCAE 為止，我的研究接續往下驗證 token 分配品質。")

# Slide 5: 老師回饋與方向調整
add_content_slide(prs, "老師回饋與方向調整", [
    ("老師的建議（7/4 討論後）", 0),
    ("  先做 pilot 實驗驗證方向可行性", 1),
    ("  用 Process Mining 分析推理軌跡——新的切入點", 1),
    ("  設計實驗收集初步數據，再決定論文方向", 1),
    ("", 0),
    ("我的調整", 0),
    ("  V1 Pilot：快速驗證 PM 能否分析 LLM 推理軌跡", 1),
    ("  → 確認 PM 能區分模型推理風格", 1),
    ("  V2：擴大實驗，加入校準指標，修正 V1 缺陷", 1),
    ("  → 發現信心差距、推理風格散度等新指標", 1),
])

# Slide 6: Section - V1 Pilot
add_section_slide(prs, "2. V1 Pilot 實驗", "GSM8K · 20 題 · 5 模型")

# Slide 7: V1 設計
add_table_slide(prs, "V1 實驗設計",
    ["參數", "設定"],
    [
        ["題目", "20 題 GSM8K（小學數學）"],
        ["模型", "5 個（GPT-OSS-20B/120B, DeepSeek, GLM-4.7, GLM-5.2）"],
        ["活動類型", "8 種（understand, recall, plan, calculate, reason, verify, reconsider, answer）"],
        ["PM 分析", "Petri net 發現 + 一致性檢查"],
        ["校準", "未收集"],
        ["準確率", "95–100%（太簡單）"],
    ],
    notes="V1 的主要問題是題目太簡單，準確率都接近滿分，沒有變異量做相關分析。")

# Slide 8: V1 結果 - 準確率與軌跡
add_two_image_slide(prs, "V1 結果：準確率與路徑結構",
    f"{FIG_V1}/fig1_accuracy_token_efficiency.png",
    f"{FIG_V1}/fig2_path_structure.png",
    "準確率與 Token 效率", "路徑結構指標")

# Slide 9: V1 發現
add_content_slide(prs, "V1 關鍵發現", [
    ("✅ PM 能區分推理風格", 0),
    ("  DeepSeek = 直覺型（短軌跡、高 answer 比例）", 1),
    ("  GPT-OSS-20B = 掙扎型（長軌跡、高變異、reason 重）", 1),
    ("  GLM-5.2 = 混合型（變體多、有驗證）", 1),
    ("", 0),
    ("❌ 限制", 0),
    ("  題目太簡單 → 準確率 95-100%，無法做相關分析", 1),
    ("  無校準數據 → 無法驗證校準與 token 的關係", 1),
    ("  變體太多 → Petri net 太寬鬆", 1),
    ("", 0),
    ("→ 需要更難的題目 + 收集信心數據", 0),
])

# Slide 10: Section - V2
add_section_slide(prs, "3. V2 實驗設計與執行", "MMLU STEM + ARC · 100 題 · 4 模型")

# Slide 11: V1 vs V2 比較
add_table_slide(prs, "V1 → V2 改動",
    ["面向", "V1", "V2"],
    [
        ["題目", "20 GSM8K", "100 (MMLU STEM 50 + ARC 50)"],
        ["難度", "簡單 (95-100%)", "中等 (56-98%)"],
        ["模型", "5", "4 (GLM-4.7 下架)"],
        ["信心自評", "無", "多輪對話式 (0-100%)"],
        ["校準指標", "無", "Brier, 信心差距"],
        ["PM 分析", "Petri net", "Petri net + 熵 + JSD"],
    ],
    notes="V2 的三個關鍵改動：更難的題目、信心收集、無參考模型分析方法。")

# Slide 12: V2 模型
add_table_slide(prs, "V2 模型選擇",
    ["模型", "參數量", "架構", "準確率", "平均 Token"],
    [
        ["GPT-OSS-20B", "21B", "Dense", "56%", "773"],
        ["DeepSeek-V4-Flash", "158B", "MoE (13B active)", "97%", "809"],
        ["GPT-OSS-120B", "117B", "Dense", "79%", "619"],
        ["GLM-5.2", "756B", "MoE (40B active)", "98%", "823"],
    ])

# Slide 13: 信心 Prompt 演進
add_content_slide(prs, "信心 Prompt 的演進（本身就是研究發現）", [
    ("V1（無上下文）", 0),
    ("  「你回答了 D，你有多確定？」", 1),
    ("  → DeepSeek 回傳 2%, GLM-5.2 回傳 27%（無效數據）", 1),
    ("", 0),
    ("V2（多輪對話）", 0),
    ("  Turn 1: 原始題目 + 選項", 1),
    ("  Turn 2: 模型完整推理回應", 1),
    ("  Turn 3: 「根據你的推理，你有多確定？」", 1),
    ("  → DeepSeek 99%, GLM-5.2 99%（有意義的數據）", 1),
    ("", 0),
    ("教訓：無上下文的信心 prompt 會產生完全失真的數據", 0),
    ("這本身就是一個 methodology finding", 0),
])

# Slide 14: Section - V2 結果
add_section_slide(prs, "4. V2 關鍵發現")

# Slide 15: 準確率與 Token
add_image_slide(prs, "準確率與 Token 效率",
    f"{FIG}/fig1_accuracy_token_efficiency.png",
    "56-98% 準確率範圍，足夠做相關分析",
    notes="V2 成功創造了準確率變異。GLM-5.2 和 DeepSeek 幾乎並列，但參數差 5.8 倍。")

# Slide 16: 路徑結構
add_image_slide(prs, "推理路徑結構",
    f"{FIG}/fig2_path_structure.png",
    "DeepSeek 最短（9.1步）、GPT-OSS-20B 最長（16.9步）",
    notes="DeepSeek 從不自我修正，GPT-OSS-20B 變異極大。")

# Slide 17: 校準指標
add_image_slide(prs, "校準分析：信心差距是關鍵指標",
    f"{FIG}/fig4_calibration.png",
    "",
    notes="核心發現：信心差距與準確率反相關。GPT-OSS-20B 雖然最不準，但最能偵測自己的錯誤。")

# Slide 18: 校準表格
add_table_slide(prs, "校準指標總覽",
    ["模型", "準確率", "Brier", "答對信心", "答錯信心", "差距"],
    [
        ["GLM-5.2", "98%", "0.011", "99.5%", "50.0%", "+49.5"],
        ["DeepSeek", "97%", "0.021", "99.6%", "66.7%", "+33.0"],
        ["GPT-OSS-20B", "56%", "0.047", "95.6%", "14.1%", "+81.4"],
        ["GPT-OSS-120B", "79%", "0.052", "94.5%", "28.8%", "+65.7"],
    ],
    notes="Brier 都小於 0.05，都算不錯。差距才是區分指標——GPT-OSS-20B 最強。")

# Slide 19: 信心差距含義
add_content_slide(prs, "核心發現：準確率 vs 信心差距反相關", [
    ("現象", 0),
    ("  最低準確率的 GPT-OSS-20B（56%）→ 信心差距最大（+81.4）", 1),
    ("  最高準確率的 DeepSeek（97%）→ 信心差距最小（+33.0）", 1),
    ("", 0),
    ("解讀", 0),
    ("  GPT-OSS-20B 常答錯 → 「知道自己不會」的經驗豐富 → 答錯時 14% 信心", 1),
    ("  DeepSeek 很少答錯 → 對不確定性缺乏「練習」→ 答錯時仍 67% 信心", 1),
    ("", 0),
    ("呼應學姐論文：能力強 ≠ 自評準", 0),
    ("", 0),
    ("實際含義", 0),
    ("  GPT-OSS-20B 可用信心閾值過濾錯誤 → 高精確度部署", 1),
    ("  DeepSeek 答錯時仍自信 → 很難透過自評偵測錯誤", 1),
])

# Slide 20: Petri net 問題
add_content_slide(prs, "Petri net 一致性檢查的局限", [
    ("問題", 0),
    ("  每個模型 ~90-100 個獨特變體（每題走法都不同）", 1),
    ("  Inductive Miner 產出接近 flower model 的網 → 什麼都接受", 1),
    ("  所有模型 alignment 偏離 = 0 → 沒有區分力", 1),
    ("", 0),
    ("測試換參考模型（GLM-5.2 → DeepSeek）", 0),
    ("  結果完全相同：alignment 全是 0", 1),
    ("  → 這不是參考選擇問題，是發現方法的限制", 1),
    ("", 0),
    ("解法：改用不需要參考模型的分析方法", 0),
    ("  A) 熵分析 + Levenshtein 距離", 1),
    ("  C) 步驟類型頻率 + Jensen-Shannon 散度", 1),
    ("  （參考 Berti et al. 2025, Back et al. 2019）", 1),
])

# Slide 21: 軌跡熵
add_image_slide(prs, "A1: 軌跡熵分析",
    f"{FIG}/fig_a1_trace_entropy.png",
    "每軌跡熵（左）差不多；變體熵（右）GLM-5.2=1.0（每題都不同）",
    notes="熵分析不需要參考模型。GLM-5.2 的變體熵=1.0，確認 Petri net 對它無意義。")

# Slide 22: Levenshtein 距離
add_two_image_slide(prs, "A3: Levenshtein 距離與分群",
    f"{FIG}/fig_a3_levenshtein_heatmap.png",
    f"{FIG}/fig_a3b_dendrogram.png",
    "模型間距離熱圖", "階層分群")

# Slide 23: 步驟類型頻率
add_image_slide(prs, "C1: 步驟類型頻率分佈",
    f"{FIG}/fig_c1_step_frequency.png",
    "DeepSeek: 18% understand + 28% answer（直覺型）；GPT-OSS-20B: 36% reason（掙扎型）",
    notes="步驟類型頻率比 Levenshtein 更能區分模型。DeepSeek 的分佈明顯不同。")

# Slide 24: JSD
add_image_slide(prs, "C2/C3: Jensen-Shannon 散度",
    f"{FIG}/fig_c2_c3_jsd_heatmaps.png",
    "左：步驟類型 JSD；右：Bigram JSD（轉換模式差異更大）",
    notes="DeepSeek vs GLM-5.2 差異最大（0.22/0.36），儘管準確率只差 1%。")

# Slide 25: JSD 核心發現
add_content_slide(prs, "核心發現：相同準確率，不同推理風格", [
    ("DeepSeek vs GLM-5.2", 0),
    ("  準確率：97% vs 98%（幾乎相同）", 1),
    ("  JSD（步驟類型）：0.22（差最多）", 1),
    ("  JSD（bigram）：0.36（差更多）", 1),
    ("", 0),
    ("DeepSeek 的風格", 0),
    ("  理解 → 直接作答（18% understand + 28% answer + 21% reason）", 1),
    ("", 0),
    ("GLM-5.2 的風格", 0),
    ("  計算 + 推理 + 驗證（34% calculate + 31% reason + 12% understand + 12% answer）", 1),
    ("", 0),
    ("→ 兩條路都到 98%，但結構完全不同", 0),
    ("→ 這只有步驟類型散度看得到，accuracy 和 Brier 都看不出來", 0),
])

# Slide 26: Top transitions
add_image_slide(prs, "各模型 Top-10 轉換模式",
    f"{FIG}/fig_c3c_top_transitions.png",
    "DeepSeek: understand→answer 為主；其他模型: reason→calculate 來回",
    notes="DeepSeek 的轉換模式最簡潔。")

# Slide 27: 四種推理風格
add_table_slide(prs, "四種推理風格（V1 + V2 一致）",
    ["風格", "模型", "特徵", "準確率", "信心差距"],
    [
        ["直覺型", "DeepSeek", "短軌跡、高 answer、無迴圈", "97%", "+33.0"],
        ["系統型", "GPT-OSS-120B", "中等、calculate+reason 均衡", "79%", "+65.7"],
        ["混合型", "GLM-5.2", "變體多、有驗證、最均衡", "98%", "+49.5"],
        ["掙扎型", "GPT-OSS-20B", "長軌跡、高變異、reason 重", "56%", "+81.4"],
    ],
    notes="四種風格跨 V1 V2 穩定出現。注意準確率與信心差距的反相關。")

# Slide 28: Section - 下一步
add_section_slide(prs, "5. 下一步")

# Slide 29: 下一步
add_content_slide(prs, "下一步規劃", [
    ("立即（1-2 週）", 0),
    ("  計算 IRT 參數 → 完整 LCAE 分數", 1),
    ("  LLM 輔助步驟分割 → 驗證規則式分割可靠性", 1),
    ("  逐題分析模型分歧案例", 1),
    ("", 0),
    ("短期（2-4 週）", 0),
    ("  擴大題目集至 200–500 題", 1),
    ("  增加模型至 6–8 個", 1),
    ("  LCAE ↔ PM 指標 ↔ Token 效率 三方相關分析", 1),
    ("  信心差距作為部署指標的正式化", 1),
    ("", 0),
    ("中期（1-3 月）", 0),
    ("  撰寫論文（IEEE Big Data 2026 或 BPM/ICPM）", 1),
    ("  醫療場域延伸（NHRI 連結）", 1),
])

# Slide 30: 給老師的問題
add_content_slide(prs, "想跟老師討論的問題", [
    ("Q1: 方向確認", 0),
    ("  PM 分析 + 校準指標 + Token 效率 的三軸框架是否可行？", 1),
    ("", 0),
    ("Q2: 論文投稿目標", 0),
    ("  IEEE Big Data 2026（偏資料分析）vs BPM/ICPM（偏流程挖掘）？", 1),
    ("", 0),
    ("Q3: 活動標註可靠性", 0),
    ("  規則式標註的已知限制，LLM 輔助驗證的優先級？", 1),
    ("", 0),
    ("Q4: 與學姐框架的銜接", 0),
    ("  目前缺 IRT 參數，是否可用學姐的模型/benchmark 資料？", 1),
    ("", 0),
    ("Q5: 醫療場域", 0),
    ("  信心差距 + 推理風格分析是否適合臨床 LLM 場景？", 1),
])

# Slide 31: Summary
add_content_slide(prs, "Summary", [
    ("V1 → V2 的進展", 0),
    ("  V1 確認 PM 能分析 LLM 推理軌跡，但題目太簡單", 1),
    ("  V2 加入校準指標，發現信心差距是關鍵", 1),
    ("", 0),
    ("V2 三個核心發現", 0),
    ("  1. 信心差距與準確率反相關 — 呼應學姐「能力強 ≠ 自評準」", 1),
    ("  2. Petri net 在高變異下失效 — 改用熵 + JSD 解決", 1),
    ("  3. 相同準確率可由完全不同的推理風格達成", 1),
    ("     DeepSeek「理解→作答」 vs GLM-5.2「計算+推理+驗證」", 1),
    ("", 0),
    ("三軸分析框架", 0),
    ("  校準（Brier + 信心差距）× 推理結構（JSD + 步驟頻率）× Token 效率", 1),
])

# Slide 32: Thanks
add_title_slide(prs, "謝謝老師", "歡迎討論與指導")

# Save
output_path = "experiments/v2-mmlu-arc/report/presentation-2026-07-17.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")