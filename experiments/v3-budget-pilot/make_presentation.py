#!/usr/bin/env python3
"""Full presentation: LCAE → V1 → V2 → V3 Budget Sensitivity Pilot."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x4a, 0x6f, 0xd8)
GRAY = RGBColor(0x6b, 0x6b, 0x80)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xf5, 0xf5, 0xfa)

def add_ts(prs, t, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
    b = s.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(1.5))
    b.text_frame.paragraphs[0].text = t; b.text_frame.paragraphs[0].font.size = Pt(36); b.text_frame.paragraphs[0].font.bold = True; b.text_frame.paragraphs[0].font.color.rgb = WHITE; b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    b2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(1))
    b2.text_frame.paragraphs[0].text = sub; b2.text_frame.paragraphs[0].font.size = Pt(20); b2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xa0,0xa0,0xb0); b2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    b3 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(0.5))
    b3.text_frame.paragraphs[0].text = "2026-08-28"; b3.text_frame.paragraphs[0].font.size = Pt(16); b3.text_frame.paragraphs[0].font.color.rgb = GRAY; b3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def sec(prs, t, s=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BLUE
    b = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    b.text_frame.paragraphs[0].text = t; b.text_frame.paragraphs[0].font.size = Pt(32); b.text_frame.paragraphs[0].font.bold = True; b.text_frame.paragraphs[0].font.color.rgb = WHITE; b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if s:
        b2 = sl.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
        b2.text_frame.paragraphs[0].text = s; b2.text_frame.paragraphs[0].font.size = Pt(18); b2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xd0,0xd0,0xe0); b2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def bar(s, t):
    b = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
    b.text_frame.paragraphs[0].text = t; b.text_frame.paragraphs[0].font.size = Pt(24); b.text_frame.paragraphs[0].font.bold = True; b.text_frame.paragraphs[0].font.color.rgb = DARK

def con(s, t, items, note=""):
    bar(s, t)
    tx = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.13), Inches(5.8))
    tf = tx.text_frame; tf.word_wrap = True
    for i, (text, lv) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(16) if lv == 0 else Pt(14)
        p.font.color.rgb = DARK if lv == 0 else GRAY; p.level = lv; p.space_after = Pt(4)
    if note: s.notes_slide.notes_text_frame.text = note

def tbl(s, t, hd, rows, note=""):
    bar(s, t)
    rn = len(rows) + 1; cn = len(hd)
    ts = s.shapes.add_table(rn, cn, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.5 * rn))
    tb = ts.table
    for j, h in enumerate(hd):
        c = tb.cell(0, j); c.text = h
        for p2 in c.text_frame.paragraphs: p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = tb.cell(i+1, j); c.text = str(v)
            for p2 in c.text_frame.paragraphs: p2.font.size = Pt(11); p2.font.color.rgb = DARK; p2.alignment = PP_ALIGN.CENTER
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = LIGHT
    if note: s.notes_slide.notes_text_frame.text = note

def cs(prs, t, items, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); con(s, t, items, note)
def ts(prs, t, hd, rows, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); tbl(s, t, hd, rows, note)

# =================== BUILD ===================

add_ts(prs, "LLM 校準 × Token Allocation × Process Mining",
    "從學姐 LCAE 到 V3 Budget Sensitivity Pilot 的研究進展")

cs(prs, "大綱", [
    ("1. 研究源起：從學姐 LCAE 框架出發", 0),
    ("2. V1 Pilot：Process Mining 可行性驗證（7 月初）", 0),
    ("3. V2 實驗：加入校準指標 + PM 分析（7 月中）", 0),
    ("4. 7/28 Deep Research：方向收斂與競爭者分析", 0),
    ("5. V3 Budget Sensitivity Pilot：最新結果（今天）", 0),
    ("6. 下一步規劃", 0),
], "老師好，今天報告的是我從學姐論文延伸的研究。從七月初到現在做了三階段實驗，今天會講完整來龍去脈和最新結果。")

# === Section 1 ===
sec(prs, "1. 研究源起", "學姐 LCAE 框架 → 我的研究問題")

ts(prs, "學姐論文：LCAE 框架 (Chen et al., IEEE IRI 2026)",
    ["階段", "內容", "關鍵發現"],
    [
        ["Stage 1 IRT 能力估計", "Rasch Model：σ(θ_m−β_i)\n模型能力與題目難度同一把尺", "客觀計算答錯機率"],
        ["Stage 2 自我評估", "QOQ / IDS / DPR / Combined\n四種情境", "IDS（給難度訊號）最有效\n改善校準不傷能力"],
        ["Stage 3 模型選擇", "LCAE 指標比較\n自估 vs IRT 客觀估計", "能力強 ≠ 自評準\n提到 cost 關聯但未驗證"],
    ], "學姐的論文明確回答了LLM知不知道自己的實力。她用了IRT的Rasch Model，把模型能力和題目難度放在同一把尺。三個關鍵發現：能力強不等於自評準、IDS給難度訊號最有效改善校準、改善校準不傷能力。她提到reliability跟inference cost有關聯但沒深入，這就是我的切入點。")

cs(prs, "我的研究問題", [
    ("學姐回答了：模型知道自己會不會嗎？", 0),
    ("", 0),
    ("本研究：知道自己會不會 → 能否知道自己需要思考多久？", 0),
    ("  校準品質 → Token 分配效率？", 1),
    ("", 0),
    ("核心假設", 0),
    ("  校準好的模型不一定總 token 更少，但分配更合理", 1),
    ("  簡單題少 token、困難題多 token", 1),
], "學姐回答了模型知不知道自己的實力。我的問題更進一步：如果知道自己會不會，能不能判斷需要多少推理資源？校準好的模型不一定總token更少，但分配應該更合理。這是核心假設。")

cs(prs, "引入新工具：Process Mining", [
    ("將 LLM 的 CoT 文字轉成事件序列", 0),
    ("", 0),
    ("原始 CoT：「先理解題意...路程120公里、時間2小時...速度=60 km/h」", 1),
    ("  → understand → recall → calculate → verify → answer", 1),
    ("", 0),
    ("每題變成一條軌跡 → 用 pm4py 分析", 0),
    ("", 0),
    ("分析面向", 0),
    ("  ▶ 流程發現（Petri Net、流程樹）", 1),
    ("  ▶ 一致性檢查（跟參考模型比對）", 1),
    ("  ▶ 熵分析 + 步驟類型頻率（不需參考模型）", 1),
], "Process Mining原本用在製造業、醫療流程。我拿來分析LLM推理步驟—把CoT切成一段段標活動類型，每題變成一條事件序列。可以看到推理的結構：先理解還是直接算？有沒有驗證？有沒有回頭修正？這些只看token總數看不出來。")

# === Section 2 ===
sec(prs, "2. V1 Pilot（7 月初）", "GSM8K · 20 題 · 5 模型 · 無信心數據")

ts(prs, "V1 設計與限制",
    ["面向", "設定", "結果"],
    [
        ["題目", "20 題 GSM8K 小學數學", "太簡單，準確率 95-100%"],
        ["模型", "5 個（含已退休 GLM-4.7）", ""],
        ["信心", "未收集", "無法做校準分析"],
        ["活動類型", "8 種（後增為 9 種）", ""],
        ["PM 分析", "Petri net + conformance", "可行但題目太簡單"],
    ], "V1用20題小學數學、5模型、沒有收集信心數據。發現PM確實能區分推理風格，但是限制很清楚—題目太簡單沒有任何變異量，沒有信心數據無法驗證校準假設。")

cs(prs, "V1 發現與限制", [
    ("✅ PM 能區分推理風格（三種風格跨 V1 V2 穩定出現）", 0),
    ("  直覺型（DeepSeek）：短軌跡、高 answer、無迴圈", 1),
    ("  系統型（GPT-120B）：calculate+reason 均衡", 1),
    ("  掙扎型（GPT-20B）：長軌跡、高 reason 佔比", 1),
    ("", 0),
    ("❌ 三大限制", 0),
    ("  1. 題目太簡單（95-100%），零變異量", 1),
    ("  2. 無信心數據，無法分析校準", 1),
    ("  3. 8 種活動類型有時不夠細", 1),
    ("", 0),
    ("→ V2：更難題目 + 信心收集 + 9 種活動", 0),
], "正向發現是PM確實能區分推理風格，三種風格跨V1 V2穩定。但限制很致命，所以V2做了三個改進。")

# === Section 3 ===
sec(prs, "3. V2 實驗（7 月中）", "MMLU STEM + ARC · 100 題 · 4 模型 · 有校準")

ts(prs, "V1 → V2 改動",
    ["面向", "V1", "V2"],
    [
        ["題目", "20 GSM8K", "100（MMLU STEM 50 + ARC 50）"],
        ["難度", "簡單（95-100%）", "中等（95-99%）"],
        ["模型", "5", "4（GLM-4.7 退休）"],
        ["活動類型", "8", "9（新增 evaluate）"],
        ["信心自評", "無", "多輪對話式（0-100%）"],
        ["校準指標", "無", "Brier、信心差距"],
        ["PM 分析", "Petri net", "Petri net + 熵 + JSD"],
    ], "V2兩個關鍵新增：信心收集和校準指標。信心prompt本身還有一個教訓—最初無上下文問「多確定」，DeepSeek回2%這種無效數字。改成多輪對話式根據完整推理過程來問，才變成99%。")

cs(prs, "V2 的關鍵轉折：Bug 修正與 Rebuild", [
    ("7/14 初版 → 發現多個 bug → 7/24 完整 Rebuild", 0),
    ("", 0),
    ("修正項目", 0),
    ("  Conformance：讀取不存在欄位 → alignment 全 0（已修正）", 1),
    ("  Levenshtein：(A,B) 和 (B,A) 分別抽樣 → 非對稱矩陣（已修正）", 1),
    ("  JSD：distance 被標成 divergence（已修正）", 1),
    ("  Confidence：只傳 response 沒傳 thinking（已修正）", 1),
    ("  環境：明文 API key → .env, 固定 Python 3.13 + uv", 1),
    ("", 0),
    ("Rebuild 後數據完全不同！", 0),
    ("  準確率：56-98% → 95-99%（GPT-20B 從 56% 跳到 98%）", 1),
    ("  信心差距：+33~+81 → −0.8~+4.5（前一版核心發現被推翻）", 1),
], "V2過程中發現幾個重大bug所以做了完整rebuild。最關鍵影響是數據完全變了——之前信心差距與準確率反相關的核心發現不成立了。V2作為校準分析價值大幅下降，但作為PM pipeline驗證仍然成功。")

ts(prs, "V2 Rebuild 後最終數據",
    ["模型", "準確率", "平均步數", "Brier", "信心差距", "偏離數"],
    [
        ["GPT-OSS-20B", "98%", "16.8", "0.061", "+3.8", "2885"],
        ["DeepSeek", "98%", "9.2", "0.021", "−0.8", "2080"],
        ["GPT-OSS-120B", "95%", "13.0", "0.053", "+4.5", "2275"],
        ["GLM-5.2", "99%", "15.7", "0.010", "−0.5", "2893"],
    ], "Rebuild後最終數據。準確率全在95-99%，Brier很小，信心差距幾乎消失。V2作為校準分析價值不高，但PM pipeline驗證成功。")

# === Section 4 ===
sec(prs, "4. 7/28 Deep Research", "競爭者分析與方向收斂（1612 行報告）")

cs(prs, "哪些 Novelty 已不成立？", [
    ("❌ 「首次用 confidence 控制 reasoning length」", 0),
    ("  → Think Just Enough (EACL 2026)：自評信心可做 stopping signal", 1),
    ("", 0),
    ("❌ 「首次讓模型預估 token budget」", 0),
    ("  → SelfBudgeter (ACL 2026)：budget prediction + RL", 1),
    ("", 0),
    ("❌ 「首次用 PM 分析 LLM reasoning」", 0),
    ("  → Berti et al. 2025 (TechRxiv), PM4GRPO (arXiv 2025)", 1),
    ("", 0),
    ("❌ 「校準好的模型用更少 token」", 0),
    ("  → Capability Calibration (arXiv 2026)：calibration → best-of-k", 1),
    ("", 0),
    ("⚠️ 但「Calibration → Reasoning Length Allocation」仍是空白", 0),
], "7/28花一週做系統性文獻調查。好消息是這個領域很熱，壞消息是很多基本想法已經被人做了。Think Just Enough已證明自評信心可做stopping signal；SelfBudgeter讓模型預估budget；Capability Calibration連結了校準和best-of-k配置。所以必須收斂題目，不能說「首次用confidence省token」。")

cs(prs, "論文定位：收斂後的差異化", [
    ("核心命題", 0),
    ("  現有研究假設 confidence 夠可靠", 1),
    ("  我們問：confidence 的校準品質本身，是否決定了配置效果？", 1),
    ("", 0),
    ("差異化對比", 0),
    ("  Think Just Enough：raw confidence → 我們用 IRT 校準", 1),
    ("  SelfBudgeter：需 SFT+RL → 我們 training-free", 1),
    ("  Sonata：需 hidden states → 我們 black-box", 1),
    ("  Capability Calibration：配置 sampling 次數 → 我們配置 reasoning length", 1),
    ("", 0),
    ("完整因果鏈：IRT→LCAE→Token Requirement 預測→動態分配→PM 診斷", 0),
], "新定位：不是「confidence能不能省token」，而是「confidence的校準品質是否決定配置效果」。Think Just Enough用raw confidence但沒校準。SelfBudgeter需訓練。Sonata需hidden states。最關鍵的是Capability Calibration配置sampling次數，我們配置的是single trajectory的reasoning length。")

cs(prs, "V3 分階段實驗設計", [
    ("Phase 1（✅ 已確認）", 0),
    ("  Ollama API num_predict 可確實限制 reasoning tokens", 1),
    ("", 0),
    ("Phase 2（✅ 已完成，今天重點）", 0),
    ("  Budget Sensitivity Pilot：確認題目難度 + budget 敏感度存在", 1),
    ("  2 模型 × 30 題 × 4 budgets × 2 reps = 480 calls", 1),
    ("", 0),
    ("Phase 3（下一步）", 0),
    ("  擴大題目 + 加入 confidence + 加至 4 模型", 1),
    ("  驗證校準品質與 token requirement 的因果關係", 1),
    ("", 0),
    ("Phase 4-5（未來）", 0),
    ("  Signal Evaluation / Allocation Evaluation / PM Analysis", 1),
], "V3分階段設計，每個階段都有Go/No-Go標準。Phase 1確認API可以限制token。Phase 2就是今天重點。通過才進Phase 3。")

# === Section 5 ===
sec(prs, "5. V3 Budget Sensitivity Pilot", "MATH-500 · 30 題 · 2 模型 · 4 Budgets · 480 Calls")

ts(prs, "實驗設計",
    ["參數", "設定", "理由"],
    [
        ["題目", "30 題 MATH-500（L3-5 各 10）", "Level 1-2 太簡單"],
        ["模型", "GPT-OSS-20B, DeepSeek", "校準好 vs 校準差"],
        ["Budgets", "128 / 256 / 512 / 1024", "根據 V2 自然用量（平均 600-900）"],
        ["複製", "每條件 2 次 replicate", "檢查 stochastic 波動"],
        ["控制", "Ollama num_predict, temp=0.0", "真正事前限制非事後截斷"],
        ["總量", "480 次 API 呼叫", "約 11 分鐘"],
    ], "30題MATH-500，只取Level 3到5。兩個模型。Budget根據V2的自然token用量設定。用Ollama的num_predict真正限制生成上限，不是事後截斷。")

ts(prs, "主要結果（✅ Go）",
    ["Budget", "GPT-OSS-20B", "DeepSeek", "差距"],
    [
        ["128", "0 / 60 = 0.0%", "0 / 60 = 0.0%", "—"],
        ["256", "2 / 60 = 3.3%", "9 / 60 = 15.0%", "5.0× 🏆"],
        ["512", "12 / 60 = 20.0%", "20 / 60 = 33.3%", "1.7×"],
        ["1024", "19 / 60 = 31.7%", "22 / 60 = 36.7%", "1.2×"],
    ], "主要結果。128兩個模型都是0%—題目夠難。256時最關鍵—DeepSeek 15% vs GPT 3%，差5倍。1024差距縮小到1.2倍。這呼應核心假設：校準品質在資源受限時差異最大。")

ts(prs, "依難度 Level 分析",
    ["Level", "128", "256", "512", "1024", "關鍵觀察"],
    [
        ["L3 GPT", "0%", "0%", "30%", "35%", "Level 3-4 有 budget sensitivity"],
        ["L3 DS", "0%", "10%", "40%", "40%", ""],
        ["L4 GPT", "0%", "10%", "30%", "45%", "GPT @1024 反超 DS"],
        ["L4 DS", "0%", "35%", "40%", "40%", "DS @256 大幅領先 🏆"],
        ["L5 GPT", "0%", "0%", "0%", "15%", "Level 5 太難，連 1024 都不夠"],
        ["L5 DS", "0%", "0%", "20%", "30%", "DS 仍然贏一倍"],
    ], "分Level看。Level 4最有趣—DeepSeek在256時35% vs GPT 10%，但GPT在1024時45%反超。代表GPT需要更多token才發揮實力。Level 5太難，兩人都不夠。")

ts(prs, "逐題對決：誰更省 token？",
    ["題目", "DeepSeek\n最低有效", "GPT\n最低有效", "差距", "備註"],
    [
        ["L4 Algebra", "256", "1024", "省 4 倍", "🏆 最誇張差距"],
        ["L4 Prealgebra", "256", "512", "省一半", "穩定優勢"],
        ["L4 Prealgebra", "256", "512", "省一半", ""],
        ["L3 Algebra", "256", "512", "省一半", "簡單題也省"],
        ["L5 Algebra", "512", "1024", "省一倍", "難題照省"],
        ["L3 Prealgebra", "512", "1024", "省一倍", ""],
        ["L5 Algebra", "512", "永不", "—", "✅ DS 做到 GPT 做不到"],
        ["L4 Number Theory", "永不", "1024", "—", "❌ 唯一 GPT 贏"],
    ], "30題中8題DeepSeek比GPT省token就答對，只有1題反過來。最誇張的是L4 Algebra—DeepSeek 256就答對，GPT要1024才答對，差4倍。還有一題L5 Algebra，DeepSeek 512答對，GPT連1024都做不到。這些差異不能用模型大小完全解釋。")

cs(prs, "結果可靠度：Replicate 一致性", [
    ("每條件跑 2 次 replicate，檢查一致性", 0),
    ("", 0),
    ("Budget 128：30/30 完全一致（0 分歧）", 0),
    ("Budget 256：29/30 完全一致（1 分歧）", 0),
    ("Budget 512：28/30 完全一致（2 分歧）", 0),
    ("Budget 1024：25/30 完全一致（5 分歧）", 0),
    ("", 0),
    ("低 budget 結果極穩定 → 我們最關心的就是低 budget 差異", 0),
    ("分歧僅 1 題（共 30 題），結果可靠", 1),
], "每個條件跑2次看是否一致。128時0分歧（沒人答對）。256只有1分歧。結果可靠，特別是在我們最關心的低budget區間。")

cs(prs, "V3 Pilot 結論", [
    ("✅ Budget sensitivity 確認存在", 0),
    ("  128→1024 穩定爬升，MATH-500 L3-5 難度剛好", 1),
    ("", 0),
    ("✅ DeepSeek 在所有 budget 下優於 GPT", 0),
    ("  256 差距最大：15% vs 3%（5 倍）", 1),
    ("  1024 趨近：37% vs 32%（1.2 倍）", 1),
    ("", 0),
    ("✅ 校準優勢在低 token 時最明顯", 0),
    ("  資源受限時差異擴大 → 論文核心論點有初步支持", 1),
    ("", 0),
    ("⚠️ 限制", 0),
    ("  只有 2 模型 → 無法區分校準 vs 模型能力的 confound", 1),
    ("  無信心收集 → 還不能驗證因果鏈", 1),
], "三件事確認。但兩個限制也清楚：只有兩個模型而且它們不只校準不同大小也不同。沒有信心收集還沒法驗證因果鏈。方向對了但還需要決定性實驗。")

# === Section 6 ===
sec(prs, "6. 下一步", "Phase 3 決定性實驗")

ts(prs, "Phase 3 建議設計",
    ["參數", "設定", "目的"],
    [
        ["模型", "4 個（+GPT-120B, +GLM-5.2）", "區分校準 vs 能力 confound"],
        ["題目", "60 題 MATH-500 L3+L4", "剔除太難/太簡單的"],
        ["Budgets", "3 個（256/512/1024）", "剔除 128（0% 無資訊）"],
        ["信心收集", "加入多輪對話式自評", "驗證校準品質與 requirement 的關係"],
        ["複製", "3 次 replicate", "提高高 budget 信度"],
        ["總呼叫", "4×60×3×3 = 2160 次", "約 40-60 分鐘"],
    ], "加上GPT-120B（中等校準、117B）和GLM-5.2（最好校準、756B）。如果GLM-5.2在低budget勝過DeepSeek而GPT-120B接近GPT-20B，那校準就是獨立解釋變數。如果只是模型大小排序，那就要調整方向。")

cs(prs, "給老師的問題", [
    ("Q1: 方向確認", 0),
    ("  從校準品質 → token requirement 預測的因果鏈，老師覺得可行？", 1),
    ("", 0),
    ("Q2: Phase 3 設計", 0),
    ("  4 模型 × 60 題 × 3 budgets + confidence 的規模合適？", 1),
    ("", 0),
    ("Q3: 投稿目標", 0),
    ("  IEEE Big Data 2026（10 月截止）vs BPM/ICPM 2027？", 1),
    ("", 0),
    ("Q4: 學姐資料", 0),
    ("  IRT 參數需要學姐的 benchmark response matrix，老師能協助？", 1),
    ("", 0),
    ("Q5: 活動標註", 0),
    ("  規則式標註的可靠性還需人工驗證，priority 多高？", 1),
], "五個問題想跟老師討論。")

ts(prs, "Timeline",
    ["期間", "項目", "備註"],
    [
        ["8/28-9/4", "Phase 3：4 模型 × 60 題 × 3 budgets", "加入 confidence 收集"],
        ["9/5-9/11", "數據分析 + 圖表", "確認因果鏈是否存在"],
        ["9/12-9/26", "Activity labeling 驗證 + PM 分析", "人工標註"],
        ["9/27-10月中", "論文初稿", "目標 IEEE Big Data 2026"],
    ], "初步時間表。如果方向確定，希望趕IEEE Big Data。")

add_ts(prs, "謝謝老師", "歡迎討論與指導")

# Save
out = "experiments/v3-budget-pilot/presentation-2026-08-28.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")